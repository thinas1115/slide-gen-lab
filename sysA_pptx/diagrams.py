"""図解系スライド: AWS構成図・ステークホルダー調整図・体制図。"""
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from generate import (ACCENT, BODY_TOP, BODY_BOTTOM, BODY_W, GRAY, LIGHT,
                      MARGIN, NAVY, TEXT, WHITE, add_rect, add_text, header,
                      note_line)
from textfit import line_height_in, text_width_in

ORANGE = RGBColor(0xE8, 0x7B, 0x1E)   # compute
GREEN = RGBColor(0x3F, 0x86, 0x24)    # storage
PURPLE = RGBColor(0x7D, 0x3F, 0x98)   # ML
LINE = RGBColor(0x59, 0x59, 0x59)


def add_arrow(slide, x1, y1, x2, y2, *, color=LINE, width=1.5, both=False,
              dash=None):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    heads = [("a:tailEnd", True), ("a:headEnd", both)]
    for tag, on in heads:
        if on:
            el = ln.makeelement(qn(tag), {"type": "triangle", "w": "med", "len": "med"})
            ln.append(el)
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.insert(0, d)
    return conn


def arrow_label(slide, cx, cy, text, w=1.6, size=9):
    """線をまたぐラベル。白背景マスクは実測テキスト幅に合わせ、wは上限として扱う
    (固定幅だと短い文字列ほど余計な範囲まで線を隠してしまうため)。
    """
    pad = 0.14
    actual_w = min(w, text_width_in(text, size) + pad)
    actual_h = line_height_in(size, 1.1) + 0.08
    tb = add_text(slide, cx - actual_w / 2, cy - actual_h / 2, actual_w, actual_h,
                  text, size, color=TEXT, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
    tb.fill.solid()
    tb.fill.fore_color.rgb = WHITE
    return tb


def node(slide, x, y, w, h, title, sub=None, *, bar=ACCENT, fill=WHITE,
         border=LINE, tsize=11, ssize=9):
    """上部にカラーバーを持つサービスノード。"""
    add_rect(slide, x, y, w, h, fill, line=border)
    add_rect(slide, x, y, w, 0.09, bar)
    ty = y + 0.12
    add_text(slide, x + 0.08, ty, w - 0.16, 0.32, title, tsize, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, x + 0.08, y + h - 0.42, w - 0.16, 0.36, sub, ssize,
                 color=GRAY, align=PP_ALIGN.CENTER)


def container(slide, x, y, w, h, label, color=LINE, dash=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.background()
    sp.line.color.rgb = color
    sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    if dash:
        ln = sp.line._get_or_add_ln()
        ln.insert(0, ln.makeelement(qn("a:prstDash"), {"val": dash}))
    add_text(slide, x + 0.12, y + 0.06, w - 0.3, 0.28, label, 10, bold=True,
             color=color)


# ---- AWS構成図(公式アイコン使用) ----
ICON_DIR = Path(__file__).parent / "assets"
ICON_R = 0.31        # アイコン半径(0.62角の半分)
EDGE_GAP = 0.06      # 矢印端点とアイコン縁の隙間


def icon_node(slide, cx, cy, img, title, sub=None, size=0.62):
    """AWS公式スタイル: アイコン+直下にサービス名ラベル。"""
    p = ICON_DIR / img
    if not p.exists():
        raise FileNotFoundError(
            f"アイコン {img} が {ICON_DIR} にありません。extract_aws_icons.py で"
            f"生成するか、ノードの icon を外して汎用図形ノードにしてください。")
    slide.shapes.add_picture(str(p), Inches(cx - size / 2),
                             Inches(cy - size / 2), Inches(size), Inches(size))
    add_text(slide, cx - 1.05, cy + size / 2 + 0.05, 2.1, 0.28, title, 10.5,
             bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, cx - 1.05, cy + size / 2 + 0.33, 2.1, 0.26, sub, 8.5,
                 color=GRAY, align=PP_ALIGN.CENTER)


def box_node(slide, cx, cy, title, sub=None, size=0.62, color=ACCENT):
    """アイコン画像を使わない汎用ノード: 角丸四角+上部カラーバー。
    icon_node と同じ外形寸法・ラベル位置なので、diagram_layout の座標計算
    (ICON_R 基準のポート・コンテナ外接)がそのまま成立する。
    アイコン素材が用意できないテーマ(オンプレNW・業務システム等)用。
    """
    sp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx - size / 2), Inches(cy - size / 2),
        Inches(size), Inches(size))
    sp.adjustments[0] = 0.12
    sp.fill.solid()
    sp.fill.fore_color.rgb = WHITE
    sp.line.color.rgb = color
    sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    add_rect(slide, cx - size / 2 + 0.08, cy - size / 2 + 0.08,
             size - 0.16, 0.07, color)
    add_text(slide, cx - 1.05, cy + size / 2 + 0.05, 2.1, 0.28, title, 10.5,
             bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, cx - 1.05, cy + size / 2 + 0.33, 2.1, 0.26, sub, 8.5,
                 color=GRAY, align=PP_ALIGN.CENTER)


def right_of(cx):
    """アイコン右縁の矢印端点X。"""
    return cx + ICON_R + EDGE_GAP


def left_of(cx):
    return cx - ICON_R - EDGE_GAP


# (旧ハンドコード版 s_aws は diagram_layout.py + diagram_specs.py に移行済み)


# ---- ステークホルダー調整図(ハブ型) ----
def s_hub(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    cx, cy = 6.67, 4.35
    hw, hh = 2.3, 1.05
    # 周辺ノード: (x, y, タイトル, 出す矢印ラベル, 戻り矢印ラベル)
    ring = spec["ring"]
    pos = [(cx - 4.6, cy - 2.15), (cx + 2.3, cy - 2.15),
           (cx - 4.6, cy + 1.25), (cx + 2.3, cy + 1.25),
           (cx - 1.15, cy - 2.55), (cx - 1.15, cy + 1.65)]
    bw, bh = 2.3, 0.85
    for (bx, by), item in zip(pos, ring):
        node(slide, bx, by, bw, bh, item["name"], item.get("sub"), bar=ACCENT)
        # 中心との接続(端点をノード側/ハブ側の縁に寄せる)
        ex = bx + bw / 2 + (0.0 if abs(bx + bw / 2 - cx) < 1 else (0.9 if bx < cx else -0.9))
        sx = bx + bw / 2
        sy = by + (bh if by < cy else 0)
        tx = cx + (-hw / 2 * 0.7 if sx < cx else hw / 2 * 0.7)
        ty = cy + (-hh / 2 if by < cy else hh / 2)
        if abs(sx - cx) < 1.2:
            tx = cx
        add_arrow(slide, sx, sy, tx, ty, both=True, color=LINE, width=1.25)
        mx, my = (sx + tx) / 2, (sy + ty) / 2
        arrow_label(slide, mx, my, item["label"], w=1.9, size=8.5)
    # 中心ハブ
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - hw / 2), Inches(cy - hh / 2),
                                Inches(hw), Inches(hh))
    sp.fill.solid()
    sp.fill.fore_color.rgb = NAVY
    sp.line.fill.background()
    sp.shadow.inherit = False
    add_text(slide, cx - hw / 2, cy - 0.3, hw, 0.6, spec["hub"], 13, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if spec.get("note"):
        note_line(slide, spec["note"])


# ---- 体制図 ----
def s_org(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])

    def box(x, y, w, h, title, sub=None, fill=WHITE, tcolor=NAVY, border=LINE):
        add_rect(slide, x, y, w, h, fill, line=border)
        add_text(slide, x + 0.06, y + (0.1 if sub else (h - 0.3) / 2), w - 0.12, 0.3,
                 title, 11.5, bold=True, color=tcolor, align=PP_ALIGN.CENTER)
        if sub:
            add_text(slide, x + 0.06, y + h - 0.38, w - 0.12, 0.32, sub, 8.5,
                     color=GRAY, align=PP_ALIGN.CENTER)

    def vline(x, y1, y2):
        add_arrow(slide, x, y1, x, y2, width=1.25)

    cx = 6.67
    # 1段目: ステコミ
    box(cx - 1.9, 1.95, 3.8, 0.62, spec["top"]["name"], None, NAVY, RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, cx + 2.1, 2.02, 3.0, 0.5, spec["top"]["sub"], 9, color=GRAY)
    # 2段目: PM
    vline(cx, 2.57, 3.05)
    box(cx - 1.9, 3.05, 3.8, 0.75, spec["pm"]["name"], spec["pm"]["sub"])
    # 幹線
    teams = spec["teams"]
    n = len(teams)
    tw, gap = 2.7, 0.5
    total = n * tw + (n - 1) * gap
    x0 = cx - total / 2
    trunk_y = 4.25
    slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(3.8),
                               Inches(cx), Inches(trunk_y)).line.color.rgb = LINE
    hl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x0 + tw / 2), Inches(trunk_y),
                                    Inches(x0 + total - tw / 2), Inches(trunk_y))
    hl.line.color.rgb = LINE
    for i, t in enumerate(teams):
        x = x0 + i * (tw + gap)
        vline(x + tw / 2, trunk_y, 4.6)
        box(x, 4.6, tw, 0.95, t["name"], t["sub"], LIGHT)
        for k, m in enumerate(t.get("members", [])):
            add_text(slide, x + 0.15, 5.62 + k * 0.3, tw - 0.3, 0.3, "・" + m, 9.5, color=TEXT)
    # 外部支援(点線)
    ex = spec["external"]
    box(10.6, 3.05, 2.15, 0.75, ex["name"], ex["sub"], WHITE)
    add_arrow(slide, cx + 1.9, 3.42, 10.6, 3.42, dash="dash", width=1.25)
    arrow_label(slide, (cx + 1.9 + 10.6) / 2, 3.24, ex["label"], w=1.7)
    if spec.get("note"):
        note_line(slide, spec["note"])
