"""生成済みPPTXのレイアウト衝突を機械検知する品質ゲート。

使い方: python check_layout.py ..\\out\\sysA_deck2.pptx
検知対象:
  T-T: テキストグリフ同士の交差
  T-P: テキストグリフ×画像(アイコン)
  T-S: テキストグリフ×塗り図形の部分重なり(完全内包は許可)
  T-F: テキストグリフがコンテナ枠線をまたぐ
  L-T: 矢印・線×テキストグリフ(白塗りマスクラベルは除外)
  L-P: 矢印・線×画像(アイコン)
  OOB: スライド境界からのはみ出し
"""
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from textfit import line_height_in, wrap_text

EMU = 914400
SLIDE_W, SLIDE_H = 13.333, 7.5
EDGE = 0.03
SEG_TRIM = 0.08
EPS = 0.03


def rect_of(sh):
    return (sh.left / EMU, sh.top / EMU, (sh.left + sh.width) / EMU,
            (sh.top + sh.height) / EMU)


def intersects(a, b, eps=EPS):
    return not (a[2] - eps <= b[0] or b[2] - eps <= a[0]
                or a[3] - eps <= b[1] or b[3] - eps <= a[1])


def contains(outer, inner, eps=EPS):
    return (outer[0] - eps <= inner[0] and outer[1] - eps <= inner[1]
            and outer[2] + eps >= inner[2] and outer[3] + eps >= inner[3])


def glyph_rect(sh):
    """add_text/セルの実グリフ矩形をtextfitで推定する。"""
    tf = sh.text_frame
    box = rect_of(sh)
    bw = box[2] - box[0]
    lines = []
    align = PP_ALIGN.LEFT
    for p in tf.paragraphs:
        if not p.runs:
            continue
        size = p.runs[0].font.size.pt if p.runs[0].font.size else 11.0
        bold = bool(p.runs[0].font.bold)
        text = "".join(r.text for r in p.runs)
        weight = "bold" if bold else "regular"
        for ln in wrap_text(text, bw, size, weight):
            lines.append((ln, size, weight))
        if p.alignment is not None:
            align = p.alignment
    if not lines:
        return None
    from textfit import text_width_in
    w = max(text_width_in(t, s, wt) for t, s, wt in lines)
    h = sum(line_height_in(s) for _, s, _ in lines)
    if align == PP_ALIGN.CENTER:
        x0 = box[0] + (bw - w) / 2
    elif align == PP_ALIGN.RIGHT:
        x0 = box[2] - w
    else:
        x0 = box[0]
    if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
        y0 = box[1] + ((box[3] - box[1]) - h) / 2
    else:
        y0 = box[1]
    return (x0, y0, x0 + w, y0 + h)


def seg_of(sh):
    """コネクタの線分端点(flipで向きを解決)。"""
    x1, y1, x2, y2 = rect_of(sh)
    el = sh._element
    flip_h = el.spPr.xfrm.get("flipH") == "1" if el.spPr.xfrm is not None else False
    flip_v = el.spPr.xfrm.get("flipV") == "1" if el.spPr.xfrm is not None else False
    ax, bx = (x2, x1) if flip_h else (x1, x2)
    ay, by = (y2, y1) if flip_v else (y1, y2)
    return (ax, ay, bx, by)


def seg_hits_rect(seg, r):
    """線分と矩形の交差(端をSEG_TRIMだけ縮めて接続点は無視)。"""
    x1, y1, x2, y2 = seg
    import math
    length = math.hypot(x2 - x1, y2 - y1)
    if length < SEG_TRIM * 2:
        return False
    ux, uy = (x2 - x1) / length, (y2 - y1) / length
    x1, y1 = x1 + ux * SEG_TRIM, y1 + uy * SEG_TRIM
    x2, y2 = x2 - ux * SEG_TRIM, y2 - uy * SEG_TRIM
    steps = max(2, int(length / 0.05))
    for i in range(steps + 1):
        t = i / steps
        px, py = x1 + (x2 - x1) * t, y1 + (y2 - y1) * t
        if r[0] + EPS < px < r[2] - EPS and r[1] + EPS < py < r[3] - EPS:
            return True
    return False


def has_solid_fill(sh):
    try:
        return sh.fill.type is not None and str(sh.fill.type) == "SOLID (1)"
    except Exception:
        return False


def snippet(t):
    t = t.replace("\n", " ")
    return t[:16] + ("…" if len(t) > 16 else "")


def check(path):
    prs = Presentation(path)
    findings = []
    for si, slide in enumerate(prs.slides, 1):
        texts, pics, solids, frames, segs = [], [], [], [], []
        for z, sh in enumerate(slide.shapes):
            st = sh.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                pics.append((rect_of(sh), sh.name))
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                (solids if has_solid_fill(sh) else frames).append(
                    (rect_of(sh), sh.name, z))
            elif st == MSO_SHAPE_TYPE.LINE:
                segs.append((seg_of(sh), sh.name, z))
            elif sh.has_text_frame and sh.text_frame.text.strip():
                g = glyph_rect(sh)
                if g:
                    texts.append((g, snippet(sh.text_frame.text), has_solid_fill(sh)))

        for i in range(len(texts)):
            for j in range(i + 1, len(texts)):
                if intersects(texts[i][0], texts[j][0]):
                    findings.append((si, "T-T", texts[i][1], texts[j][1]))

        for g, t, _ in texts:
            for r, name in pics:
                if intersects(g, r):
                    findings.append((si, "T-P", t, name))

        for g, t, _ in texts:
            for r, name, _z in solids:
                if intersects(g, r) and not contains(r, g):
                    findings.append((si, "T-S", t, name))

        for g, t, masked in texts:
            if masked:
                continue
            for r, name, _z in frames:
                inner = (r[0] + EDGE, r[1] + EDGE, r[2] - EDGE, r[3] - EDGE)
                if intersects(g, r, eps=-EDGE) and not contains(inner, g) \
                        and intersects(g, r):
                    findings.append((si, "T-F", t, name))

        for seg, name, lz in segs:
            for g, t, masked in texts:
                if masked or not seg_hits_rect(seg, g):
                    continue
                covered = any(contains(r, g) and sz > lz for r, _n, sz in solids)
                if not covered:
                    findings.append((si, "L-T", name, t))

        for seg, name, _lz in segs:
            for r, pname in pics:
                if seg_hits_rect(seg, r):
                    findings.append((si, "L-P", name, pname))

        for g, t, _ in texts:
            if g[0] < 0 or g[1] < 0 or g[2] > SLIDE_W or g[3] > SLIDE_H:
                findings.append((si, "OOB", t, ""))
    return findings


if __name__ == "__main__":
    fs = check(sys.argv[1])
    if not fs:
        print(f"OK: no layout collisions in {sys.argv[1]}")
        sys.exit(0)
    print(f"NG: {len(fs)} finding(s) in {sys.argv[1]}")
    for si, kind, a, b in fs:
        print(f"  slide{si:02d} [{kind}] {a!r} x {b!r}")
    sys.exit(1)
