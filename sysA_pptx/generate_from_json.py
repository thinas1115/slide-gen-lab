"""content.json を入力にしてPPTXを生成する。

使い方:
  python sysA_pptx/generate_from_json.py content.json out/from_json.pptx

生成前に validate_content.py の検証を通す。エラーがあればレンダリングせず、
生成AIにそのまま渡して直させられる粒度のメッセージを出して終了する。
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import generate
from diagrams import s_hub, s_org
from diagrams2 import s_matrix, s_process, s_roadmap
from validate_content import validate

# aws / aws2 は意図的に登録しない: 図の中身がコード内固定のサンプル専用
# rendererで、content.json からは差し替えられない。登録したままにすると
# 「タイトルだけ新規テーマ、中身は既存サンプルのAWS構成図」の資料が
# エラーなく生成されてしまう(実際に別環境の生成AIで発生した事故)。
# サンプルデッキの再生成は generate2.py / generate_patterns.py を使う。
RENDER = dict(generate.RENDER,
              hub=s_hub, org=s_org,
              process=s_process, roadmap=s_roadmap, matrix=s_matrix)


def main(json_path, out_path):
    deck = json.loads(Path(json_path).read_text(encoding="utf-8"))
    errors = validate(deck)
    if errors:
        print(f"NG: {json_path} に {len(errors)} 件の問題 (生成を中止):",
              file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        raise SystemExit(1)
    generate.DECK = deck

    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(deck["slides"], 1):
        slide = prs.slides.add_slide(blank)
        try:
            renderer = RENDER[spec["type"]]
        except KeyError as e:
            known = ", ".join(sorted(RENDER))
            raise SystemExit(f"unknown slide type: {spec['type']!r}. known: {known}") from e
        renderer(slide, spec, idx)
        if spec["type"] != "title":
            generate.footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(deck['slides'])} slides)")


if __name__ == "__main__":
    main(
        sys.argv[1] if len(sys.argv) > 1 else "content.json",
        sys.argv[2] if len(sys.argv) > 2 else "out/from_json.pptx",
    )
