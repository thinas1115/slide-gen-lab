"""社内テンプレート用のパターン検証ギャラリーを生成する。"""
import sys

from pptx import Presentation
from pptx.util import Inches

import generate
from content_patterns import PATTERN_DECK
from diagrams import s_hub, s_org
from diagrams2 import s_matrix, s_process, s_roadmap
from diagram_layout import render_diagram
from diagram_specs import DIAGRAMS


def s_diagram(slide, spec, page):
    """構成図(宣言的エンジン)。ギャラリーではインライン仕様("diagram")と
    サンプル仕様の名前参照("spec"、diagram_specs.py)の両方を許す。
    ※新規資料用の generate_from_json.py はインラインのみ(サンプル流用防止)。
    """
    generate.header(slide, spec["kicker"], spec["title"])
    d = spec["diagram"] if "diagram" in spec else DIAGRAMS[spec["spec"]]
    render_diagram(slide, d, note=spec.get("note"))


RENDER = dict(generate.RENDER,
              hub=s_hub, org=s_org,
              process=s_process, roadmap=s_roadmap, matrix=s_matrix,
              diagram=s_diagram)


def main(out_path):
    # footer() は generate.DECK の meta を参照するため、ギャラリー用に差し替える。
    generate.DECK = PATTERN_DECK
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(PATTERN_DECK["slides"], 1):
        slide = prs.slides.add_slide(blank)
        RENDER[spec["type"]](slide, spec, idx)
        if spec["type"] != "title":
            generate.footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(PATTERN_DECK['slides'])} slides)")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "../out/pattern_gallery.pptx")
