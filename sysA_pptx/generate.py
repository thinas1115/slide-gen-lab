"""システムA: python-pptx + テキスト実測によるスライド生成。"""
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from content import DECK
from textfit import fit_font_size, line_height_in, wrap_text

# ---- テーマ ----
NAVY = RGBColor(0x1F, 0x38, 0x64)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
LIGHT = RGBColor(0xEA, 0xF1, 0xF8)
TEXT = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x7F, 0x7F, 0x7F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ZEBRA = RGBColor(0xF2, 0xF6, 0xFB)
FONT = "Yu Gothic"
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.55
BODY_W = SLIDE_W - MARGIN * 2
BODY_TOP, BODY_BOTTOM = 1.62, 6.85


def set_run(run, size, *, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "ja-JP")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def add_text(slide, x, y, w, h, text, size, *, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.3,
             wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        set_run(p.add_run(), size, bold=bold, color=color)
        p.runs[0].text = line
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None, round_=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def header(slide, kicker, title):
    add_text(slide, MARGIN, 0.34, BODY_W, 0.3, kicker, 11, bold=True, color=ACCENT)
    size, _ = fit_font_size(title, BODY_W, 0.9, 20, min_pt=15, weight="bold", spacing=1.2)
    add_text(slide, MARGIN, 0.62, BODY_W, 0.9, title, size, bold=True, color=NAVY, spacing=1.2)
    add_rect(slide, MARGIN, 1.48, BODY_W, 0.022, ACCENT)


def footer(slide, page):
    add_rect(slide, MARGIN, 7.06, BODY_W, 0.012, RGBColor(0xD9, 0xD9, 0xD9))
    add_text(slide, MARGIN, 7.12, 8, 0.25, DECK["meta"]["footer"], 8, color=GRAY)
    add_text(slide, SLIDE_W - MARGIN - 1.2, 7.12, 1.2, 0.25, f"{page}", 8,
             color=GRAY, align=PP_ALIGN.RIGHT)


def note_line(slide, note):
    add_text(slide, MARGIN, 6.62, BODY_W, 0.25, note, 8.5, color=GRAY, align=PP_ALIGN.RIGHT)


# ---- スライド種別 ----
def s_title(slide, spec, page):
    add_rect(slide, 0, 0, SLIDE_W, 0.18, NAVY)
    add_rect(slide, MARGIN, 2.52, 0.7, 0.045, ACCENT)
    add_text(slide, MARGIN, 2.78, BODY_W, 1.0, spec["title"], 30, bold=True, color=NAVY)
    add_text(slide, MARGIN, 3.72, BODY_W, 0.6, spec["subtitle"], 15, color=GRAY)
    meta = DECK["meta"]
    add_text(slide, MARGIN, 6.35, BODY_W, 0.3, f"{meta['date']}    {meta['author']}", 11, color=GRAY)
    add_rect(slide, 0, 7.32, SLIDE_W, 0.18, NAVY)


def s_bullets(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    bullets = spec["bullets"]
    area_h = BODY_BOTTOM - BODY_TOP - 0.1
    tx = MARGIN + 0.42
    tw = BODY_W - 0.42
    # 内容の自然高さ+一定gapで詰め、収まる最大サイズを探す
    size, gap = 15.0, 0.42
    while size > 11:
        heights = [len(wrap_text(t, tw, size)) * line_height_in(size)
                   for t, _ in bullets]
        total = sum(heights) + gap * (len(bullets) - 1)
        if total <= area_h:
            break
        size -= 0.5
    y = BODY_TOP + 0.1 + max(0.0, (area_h - total) * 0.35)  # やや上寄せの縦バランス
    for (text, _), bh in zip(bullets, heights):
        lh = line_height_in(size)
        add_rect(slide, MARGIN + 0.05, y + lh / 2 - 0.055, 0.13, 0.13, ACCENT)
        add_text(slide, tx, y, tw, bh + 0.1, text, size)
        y += bh + gap


def s_cards(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    cards = spec["cards"]
    n = len(cards)
    gap = 0.3
    cw = (BODY_W - gap * (n - 1)) / n
    max_ch = BODY_BOTTOM - BODY_TOP - 0.15
    body_size = min(
        fit_font_size(body, cw - 0.5, max_ch - 1.05, 13, min_pt=10.5)[0]
        for _, body in cards)
    # カード高さは最長本文に合わせる(間延び防止)。縦は body 領域の中央へ
    body_h = max(len(wrap_text(b, cw - 0.5, body_size)) * line_height_in(body_size)
                 for _, b in cards)
    ch = min(max_ch, 1.0 + body_h + 0.3)
    top = BODY_TOP + 0.1 + (max_ch - ch) * 0.4
    for i, (head, body) in enumerate(cards):
        x = MARGIN + i * (cw + gap)
        add_rect(slide, x, top, cw, ch, LIGHT, round_=True)
        add_rect(slide, x + 0.25, top + 0.22, 0.32, 0.045, ACCENT)
        add_text(slide, x + 0.25, top + 0.36, cw - 0.5, 0.4, head, 14.5,
                 bold=True, color=NAVY)
        add_text(slide, x + 0.25, top + 0.9, cw - 0.5, ch - 1.05, body, body_size)


def s_table(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    cols, rows = spec["columns"], spec["rows"]
    widths = spec["col_widths"]
    assert abs(sum(widths) - BODY_W) < 0.6, f"列幅合計={sum(widths)}"
    size = 11.0
    pad = 0.09
    hdr_h = 0.38
    avail = BODY_BOTTOM - BODY_TOP - 0.15 - hdr_h - (0.3 if spec.get("note") else 0)
    while size >= 8.5:
        row_hs = []
        for row in rows:
            need = max(
                fit_font_size(c, widths[j] - pad * 2, 10, size, min_pt=size)[1].__len__()
                * line_height_in(size, 1.15) for j, c in enumerate(row))
            row_hs.append(max(need + pad * 2, 0.32))
        if sum(row_hs) <= avail:
            break
        size -= 0.5
    table_h = hdr_h + sum(row_hs)
    top = BODY_TOP + 0.1 + max(0.0, (avail - table_h) * 0.35)
    gt = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(MARGIN),
                                Inches(top), Inches(BODY_W), Inches(table_h))
    table = gt.table
    table.first_row = False
    table.horz_banding = False
    for j, wdt in enumerate(widths):
        table.columns[j].width = Emu(int(Inches(wdt)))
    table.rows[0].height = Emu(int(Inches(hdr_h)))
    for i, rh in enumerate(row_hs):
        table.rows[i + 1].height = Emu(int(Inches(rh)))
    for j, name in enumerate(cols):
        _cell(table.cell(0, j), name, size, bold=True, color=WHITE, fill=NAVY,
              center=j != len(cols) - 1 and j != 0)
    for i, row in enumerate(rows):
        fill = ZEBRA if i % 2 else WHITE
        for j, val in enumerate(row):
            _cell(table.cell(i + 1, j), val, size, fill=fill,
                  center=0 < j < len(cols) - 1 and len(val) <= 6)
    if spec.get("note"):
        note_line(slide, spec["note"])


def _cell(cell, text, size, *, bold=False, color=TEXT, fill=WHITE, center=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.09)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        p.line_spacing = 1.15
        set_run(p.add_run(), size, bold=bold, color=color)
        p.runs[0].text = line


def s_twocol(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    gap = 0.35
    cw = (BODY_W - gap) / 2
    max_ch = BODY_BOTTOM - BODY_TOP - 0.15
    panels = [spec["left"], spec["right"]]
    tw = cw - 0.68
    # 全パネル共通サイズ: 内容自然高さ+gapが収まる最大を探索
    size, bgap = 12.5, 0.26
    while size > 10:
        cont = [sum(len(wrap_text(b, tw, size)) * line_height_in(size) + bgap
                    for b in p["bullets"]) - bgap for p in panels]
        if max(cont) <= max_ch - 0.68 - 0.4:
            break
        size -= 0.5
    body_h = max(cont) + 0.44
    for i, p in enumerate(panels):
        x = MARGIN + i * (cw + gap)
        add_rect(slide, x, BODY_TOP + 0.1, cw, 0.5, NAVY, round_=True)
        add_text(slide, x + 0.22, BODY_TOP + 0.1, cw - 0.44, 0.5, p["heading"], 13.5,
                 bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, x, BODY_TOP + 0.68, cw, body_h, LIGHT, round_=True)
        y = BODY_TOP + 0.68 + 0.22
        for b in p["bullets"]:
            bh = len(wrap_text(b, tw, size)) * line_height_in(size)
            lh = line_height_in(size)
            add_rect(slide, x + 0.22, y + lh / 2 - 0.045, 0.1, 0.1, ACCENT)
            add_text(slide, x + 0.44, y, tw, bh + 0.1, b, size)
            y += bh + bgap


def s_chart(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    cd = CategoryChartData()
    cd.categories = spec["chart"]["categories"]
    for name, vals in spec["chart"]["series"]:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(MARGIN + 0.3), Inches(BODY_TOP + 0.1),
        Inches(BODY_W - 0.6), Inches(BODY_BOTTOM - BODY_TOP - 0.5), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.font.size = Pt(11)
    ch.font.name = FONT
    plot = ch.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(10.5)
    plot.gap_width = 120
    for s, colr in zip(ch.series, (RGBColor(0xBF, 0xBF, 0xBF), ACCENT)):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = colr
    if spec.get("note"):
        note_line(slide, spec["note"])


RENDER = {"title": s_title, "bullets": s_bullets, "cards": s_cards,
          "table": s_table, "twocol": s_twocol, "chart": s_chart}


def main(out_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    total = len(DECK["slides"])
    for idx, spec in enumerate(DECK["slides"], 1):
        slide = prs.slides.add_slide(blank)
        RENDER[spec["type"]](slide, spec, idx)
        if spec["type"] != "title":
            footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({total} slides)")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "../out/sysA_deck.pptx")
