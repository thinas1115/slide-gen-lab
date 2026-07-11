"""v2デッキ: 基本9枚 + 表現力検証6枚 = 15枚を生成する。

初版(sysA_deck.pptx)はそのまま残し、こちらは sysA_deck2.pptx に出力する。
"""
import sys

from pptx import Presentation
from pptx.util import Inches

import generate
from content import DECK
from content_ext import EXTRA_SLIDES
from diagrams import s_aws, s_hub, s_org
from diagrams2 import s_matrix, s_process, s_roadmap
from diagrams3 import s_aws2

RENDER2 = dict(generate.RENDER,
               aws=s_aws, hub=s_hub, org=s_org,
               process=s_process, roadmap=s_roadmap, matrix=s_matrix,
               aws2=s_aws2)


def main(out_path):
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    slides = DECK["slides"] + EXTRA_SLIDES
    for idx, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        RENDER2[spec["type"]](slide, spec, idx)
        if spec["type"] != "title":
            generate.footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(slides)} slides)")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "../out/sysA_deck2.pptx")
