"""図解系スライド第2弾: チェブロンフロー・ロードマップ・2軸マップ。"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from generate import (ACCENT, BODY_TOP, BODY_BOTTOM, BODY_W, GRAY, LIGHT,
                      MARGIN, NAVY, TEXT, WHITE, ZEBRA, add_rect, add_text,
                      header, note_line)
from diagrams import LINE, add_arrow, arrow_label
from textfit import fit_font_size

MID = RGBColor(0x9D, 0xC3, 0xE6)


# ---- チェブロン型プロセスフロー ----
def s_process(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    steps = spec["steps"]
    n = len(steps)
    gap = 0.12
    w = (BODY_W - gap * (n - 1)) / n
    y, h = 2.3, 1.05
    for i, st in enumerate(steps):
        x = MARGIN + i * (w + gap)
        shape = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.adjustments[0] = 0.28
        sp.fill.solid()
        sp.fill.fore_color.rgb = NAVY if i in spec.get("emph", []) else ACCENT
        sp.line.fill.background()
        sp.shadow.inherit = False
        add_text(slide, x + w * 0.18, y + 0.14, w * 0.72, 0.35, f"STEP {i + 1}", 9,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, x + w * 0.14, y + 0.42, w * 0.76, 0.5, st["name"], 12.5,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 下段: 説明ボックス
        size, _ = fit_font_size(st["desc"], w - 0.3, 1.55, 10.5, min_pt=9)
        add_rect(slide, x, y + h + 0.25, w, 1.9, ZEBRA, round_=True)
        add_text(slide, x + 0.15, y + h + 0.42, w - 0.3, 1.6, st["desc"], size)
        add_text(slide, x + 0.15, y + h + 1.75, w - 0.3, 0.3, st["actor"], 9,
                 bold=True, color=ACCENT)
    if spec.get("note"):
        note_line(slide, spec["note"])


# ---- ロードマップ(ガントライト) ----
def s_roadmap(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    months = spec["months"]
    label_w = 2.6
    grid_x = MARGIN + label_w
    grid_w = BODY_W - label_w
    mw = grid_w / len(months)
    top, hdr_h = BODY_TOP + 0.55, 0.42
    rows = spec["phases"]
    row_h = 1.15
    grid_h = hdr_h + len(rows) * row_h
    # 月ヘッダー
    for j, m in enumerate(months):
        add_rect(slide, grid_x + j * mw, top, mw - 0.02, hdr_h, NAVY)
        add_text(slide, grid_x + j * mw, top + 0.07, mw - 0.02, 0.3, m, 10.5,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 縦グリッド線はバー・ラベルより先に描いて背面に置く。
    for j in range(1, len(months)):
        ln = add_arrow(slide, grid_x + j * mw, top + hdr_h,
                       grid_x + j * mw, top + grid_h, width=0.5)
        ln.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
        ln.line._get_or_add_ln().remove(ln.line._get_or_add_ln().find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd"))
    # 行と帯
    for i, ph in enumerate(rows):
        ry = top + hdr_h + i * row_h
        add_rect(slide, MARGIN, ry, label_w - 0.1, row_h - 0.12,
                 LIGHT if i % 2 == 0 else ZEBRA, round_=True)
        add_text(slide, MARGIN + 0.15, ry + 0.16, label_w - 0.4, 0.35, ph["name"],
                 11.5, bold=True, color=NAVY)
        add_text(slide, MARGIN + 0.15, ry + 0.54, label_w - 0.4, 0.35, ph["goal"],
                 8.5, color=GRAY)
        x1 = grid_x + ph["start"] * mw + 0.06
        x2 = grid_x + ph["end"] * mw - 0.06
        add_rect(slide, x1, ry + 0.26, x2 - x1, 0.46, ACCENT if i % 2 == 0 else NAVY,
                 round_=True)
        add_text(slide, x1 + 0.12, ry + 0.32, x2 - x1 - 0.24, 0.32, ph["bar"], 9.5,
                 bold=True, color=WHITE)
    # マイルストーン(菱形)
    for ms in spec["milestones"]:
        mx = grid_x + ms["at"] * mw
        my = top + hdr_h + ms["row"] * row_h + 0.49
        d = 0.17
        sp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(mx - d / 2),
                                    Inches(my - d / 2), Inches(d), Inches(d))
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor(0xC0, 0x50, 0x4D)
        sp.line.fill.background()
        sp.shadow.inherit = False
        lcx = min(mx, MARGIN + BODY_W - 0.9)  # 右マージンをはみ出さない
        arrow_label(slide, lcx, my + 0.4, ms["label"], w=1.8, size=8.5)
    if spec.get("note"):
        note_line(slide, spec["note"])


# ---- 2軸ポジショニングマップ ----
def s_matrix(slide, spec, page):
    header(slide, spec["kicker"], spec["title"])
    ox, oy = 3.6, 6.35          # 原点(左下)
    aw, ah = 6.2, 4.15          # 軸の長さ
    # 象限背景(右上を強調)
    add_rect(slide, ox + aw / 2, oy - ah, aw / 2, ah / 2, LIGHT)
    # 軸
    add_arrow(slide, ox, oy, ox + aw, oy, width=1.75)
    add_arrow(slide, ox, oy, ox, oy - ah, width=1.75)
    add_text(slide, ox + aw - 2.5, oy - 0.44, 2.4, 0.3, spec["x_axis"], 10.5,
             bold=True, color=TEXT, align=PP_ALIGN.RIGHT)
    add_text(slide, ox - 0.42, oy - ah - 0.38, 3.0, 0.3, spec["y_axis"], 10.5,
             bold=True, color=TEXT)
    add_text(slide, ox + aw / 2 + 0.14, oy - ah + 0.1, 1.6, 0.3,
             spec["target_label"], 9.5, bold=True, color=ACCENT)
    # プロット
    for p in spec["points"]:
        px = ox + p["x"] * aw
        py = oy - p["y"] * ah
        r = 0.13
        sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - r), Inches(py - r),
                                    Inches(2 * r), Inches(2 * r))
        sp.fill.solid()
        sp.fill.fore_color.rgb = NAVY if p.get("emph") else GRAY
        sp.line.color.rgb = WHITE
        sp.line.width = Pt(1.0)
        sp.shadow.inherit = False
        dx, dy = p.get("lx", 0.0), p.get("ly", -0.36)
        add_text(slide, px - 0.9 + dx, py + dy, 1.8, 0.28, p["name"], 9.5,
                 bold=bool(p.get("emph")), color=NAVY if p.get("emph") else TEXT,
                 align=PP_ALIGN.CENTER)
    if spec.get("note"):
        note_line(slide, spec["note"])
