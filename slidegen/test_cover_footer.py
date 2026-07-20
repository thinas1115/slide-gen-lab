"""表紙・フッター設定の単体検証。"""
import json
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from cover_footer import (
    COVER_BACKGROUND_NAME,
    load_cover_footer_config,
    parse_cover_footer_config,
    render_cover,
    render_footer,
)
from validate_content import validate


def _must_fail(data, expected, **kwargs):
    try:
        parse_cover_footer_config(data, **kwargs)
    except ValueError as e:
        assert expected in str(e), str(e)
    else:
        raise AssertionError(f"ValueErrorになりませんでした: {data}")


def main():
    default = parse_cover_footer_config({})
    assert default.cover.eyebrow == ""
    assert default.cover.show_date is False
    assert default.cover.show_author is False
    assert default.cover.show_rail is True
    assert [item.label for item in default.cover.rail] == [
        "DATE", "ORGANIZATION", "AUTHOR",
    ]
    assert [item.value for item in default.cover.rail] == [
        "{date}", "{organization}", "{author}",
    ]
    repo_root = Path(__file__).resolve().parents[1]
    default_background = (
        repo_root / "slidegen" / "assets" / "cover" /
        "cover-background.png"
    ).resolve()
    assert default.cover.background_image == default_background
    with Image.open(default.cover.background_image) as image:
        assert image.format == "PNG"
        assert image.width / image.height > 1.7
    assert default.footer.text == "{footer}"
    assert default.footer.show_total is True

    fixture_prs = Presentation()
    fixture_prs.slide_width = Inches(13.333)
    fixture_prs.slide_height = Inches(7.5)
    fixture_slide = fixture_prs.slides.add_slide(fixture_prs.slide_layouts[6])

    example = load_cover_footer_config(repo_root / "examples" / "cover_footer.json")
    assert example.cover.background_image == default_background
    assert example.cover.eyebrow == ""
    assert example.cover.show_date is False
    assert example.cover.show_author is False
    assert example.cover.show_rail is True
    assert [item.label for item in example.cover.rail] == [
        "DATE", "ORGANIZATION", "AUTHOR",
    ]
    assert example.footer.text == "{footer}"

    solid_background = parse_cover_footer_config({
        "cover": {"background_image": None},
    })
    assert solid_background.cover.background_image is None

    minimal_deck = {
        "meta": {"title": "最小資料"},
        "slides": [{"type": "title", "title": "表紙", "subtitle": "概要"}],
    }
    assert not validate(minimal_deck)
    organization_deck = {
        "meta": {"title": "組織資料", "organization": "企画本部"},
        "slides": [{"type": "title", "title": "表紙", "subtitle": "概要"}],
    }
    assert not validate(organization_deck)
    placeholder_deck = {
        "meta": {"title": "<資料名>"},
        "slides": [{"type": "title", "title": "表紙", "subtitle": "概要"}],
    }
    assert any("入力欄が残っています" in error
               for error in validate(placeholder_deck))
    optional_text_calls = []
    render_cover(
        fixture_slide, minimal_deck["slides"][0], minimal_deck["meta"], 1,
        default,
        add_text=lambda *args, **kwargs: optional_text_calls.append(args),
        add_rect=lambda *args, **kwargs: None,
    )
    assert fixture_slide.shapes[0].name == COVER_BACKGROUND_NAME
    rendered = {args[5] for args in optional_text_calls}
    assert rendered == {"表紙", "概要"}

    default_rail_text_calls = []
    default_rail_rect_calls = []
    full_meta = {
        "title": "提案資料", "footer": "提案資料", "date": "2026年7月",
        "organization": "業務企画本部\n業務改善推進部",
        "author": "企画支援チーム",
    }
    render_cover(
        fixture_slide, {"title": "標準表紙", "subtitle": "概要"},
        full_meta, 10, default,
        add_text=lambda *args, **kwargs: default_rail_text_calls.append(args),
        add_rect=lambda *args, **kwargs: default_rail_rect_calls.append(args),
    )
    by_text = {args[5]: args for args in default_rail_text_calls}
    for text in (
        "DATE", "2026年7月", "ORGANIZATION",
        "業務企画本部\n業務改善推進部", "AUTHOR", "企画支援チーム",
    ):
        assert text in by_text
        assert by_text[text][1] == 9.82
    assert not any(
        args[5] == "2026年7月" and args[2] == 0.68
        for args in default_rail_text_calls
    )
    assert not any(
        args[5] == "企画支援チーム" and args[1] == 0.9
        for args in default_rail_text_calls
    )
    assert any(args[3] == 0.012 for args in default_rail_rect_calls)

    balanced_title_calls = []
    render_cover(
        fixture_slide,
        {"title": "プライベート接続でハマった落とし穴\n3選",
         "subtitle": "設計と運用の要点"},
        minimal_deck["meta"], 1, default,
        add_text=lambda *args, **kwargs: balanced_title_calls.append(args),
        add_rect=lambda *args, **kwargs: None,
    )
    rendered_title = next(
        args[5] for args in balanced_title_calls if "落とし穴" in args[5])
    assert rendered_title.splitlines()[-1] != "3選", rendered_title

    empty_dynamic_rail = parse_cover_footer_config({"cover": {
        "show_date": False,
        "show_author": False,
        "show_rail": True,
        "rail": [
            {"label": "DATE", "value": "{date}"},
            {"label": "ORGANIZATION", "value": "{organization}"},
            {"label": "AUTHOR", "value": "{author}"},
        ],
    }})
    empty_rail_text_calls = []
    render_cover(
        fixture_slide, minimal_deck["slides"][0], minimal_deck["meta"], 1,
        empty_dynamic_rail,
        add_text=lambda *args, **kwargs: empty_rail_text_calls.append(args),
        add_rect=lambda *args, **kwargs: None,
    )
    assert {args[5] for args in empty_rail_text_calls} == {"表紙", "概要"}

    optional_footer_calls = []
    render_footer(
        None, 1, minimal_deck["meta"], 1, default,
        add_text=lambda *args, **kwargs: optional_footer_calls.append(args),
        add_rect=lambda *args, **kwargs: None,
    )
    assert all(args[5] != "" for args in optional_footer_calls)

    custom = parse_cover_footer_config({
        "cover": {
            "eyebrow": "{date} BUSINESS REVIEW",
            "show_rail": False,
            "background_color": "102030",
        },
        "footer": {"text": "{author} | {title}", "show_total": False},
    })
    assert custom.cover.show_rail is False
    assert custom.cover.background_color == RGBColor.from_string("102030")
    assert custom.cover.title_color == default.cover.title_color
    assert custom.footer.show_total is False

    _must_fail({"cover": {"unknown": True}}, "未対応の項目")
    _must_fail({"cover": {"background_color": "blue"}}, "6桁のHEX色")
    _must_fail({"footer": {"text": "{company}"}}, "{company} は使用できません")
    _must_fail({"cover": {"rail": [{"label": "A", "value": "B"}] * 4}},
               "0〜3件")
    _must_fail({"cover": {
        "show_date": True,
        "show_rail": True,
        "rail": [{"label": "DATE", "value": "{date}"}],
    }}, "show_date を false")
    _must_fail({"cover": {
        "show_author": True,
        "show_rail": True,
        "rail": [{"label": "AUTHOR", "value": "{author}"}],
    }}, "show_author を false")

    multiline = parse_cover_footer_config({"cover": {
        "show_date": False,
        "show_author": False,
        "show_rail": True,
        "rail": [
            {"label": "DATE", "value": "{date}"},
            {"label": "ORGANIZATION",
             "value": "組織A\n部門B\nチームC"},
            {"label": "AUTHOR", "value": "担当A\n役割B"},
        ],
    }})
    text_calls = []
    rect_calls = []
    render_cover(
        fixture_slide, {"title": "Title", "subtitle": "Subtitle"},
        {"title": "T", "footer": "F", "date": "2026年7月",
         "author": "A"},
        10, multiline,
        add_text=lambda *args, **kwargs: text_calls.append((args, kwargs)),
        add_rect=lambda *args, **kwargs: rect_calls.append((args, kwargs)),
    )
    rendered_text = {args[5]: args for args, _kwargs in text_calls}
    organization = "組織A\n部門B\nチームC"
    owner = "担当A\n役割B"
    assert organization in rendered_text
    assert owner in rendered_text
    assert rendered_text[organization][4] > rendered_text[owner][4] > 0.30
    rail_rect = next(args for args, _kwargs in rect_calls if args[3] == 0.012)
    assert rail_rect[4] > 3.5

    too_many_owner_lines = parse_cover_footer_config({"cover": {
        "show_author": False,
        "show_rail": True,
        "rail": [{
            "label": "AUTHOR",
            "value": "1行目\n2行目\n3行目\n4行目",
        }],
    }})
    try:
        render_cover(
            fixture_slide, {"title": "Title", "subtitle": "Subtitle"},
            {"title": "T", "footer": "F", "date": "D", "author": "A"},
            10, too_many_owner_lines,
            add_text=lambda *args, **kwargs: None,
            add_rect=lambda *args, **kwargs: None,
        )
    except ValueError as e:
        assert "3行以内" in str(e), str(e)
    else:
        raise AssertionError("AUTHORの4行入力を拒否しませんでした")

    legacy = parse_cover_footer_config({"cover": {
        "show_date": True,
        "show_author": True,
        "show_rail": False,
    }})
    legacy_calls = []
    render_cover(
        fixture_slide, {"title": "Title", "subtitle": "Subtitle"},
        full_meta, 10,
        legacy,
        add_text=lambda *args, **kwargs: legacy_calls.append(args),
        add_rect=lambda *args, **kwargs: None,
    )
    legacy_by_text = {args[5]: args for args in legacy_calls}
    assert legacy_by_text["2026年7月"][1:3] == (9.78, 0.68)
    assert legacy_by_text["企画支援チーム"][1:3] == (0.9, 6.5)

    too_long = parse_cover_footer_config({
        "footer": {"text": "W" * 100},
    })
    try:
        render_footer(
            None, 2,
            {"title": "T", "footer": "F", "date": "D", "author": "A"},
            10, too_long,
            add_text=lambda *args, **kwargs: None,
            add_rect=lambda *args, **kwargs: None,
        )
    except ValueError as e:
        assert "footer.text の展開後の文字列が長すぎます" in str(e), str(e)
    else:
        raise AssertionError("長すぎるフッターを拒否しませんでした")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        image_path = tmp_path / "assets" / "cover.png"
        image_path.parent.mkdir()
        Image.new("RGB", (2000, 900), "#16324F").save(image_path)
        path = tmp_path / "cover_footer.json"
        path.write_text(json.dumps({
            "cover": {"background_image": "assets/cover.png"},
            "footer": {"show_text": False},
        }),
                        encoding="utf-8")
        loaded = load_cover_footer_config(path)
        assert loaded.footer.show_text is False
        assert loaded.cover.eyebrow == default.cover.eyebrow
        assert loaded.cover.background_image == image_path.resolve()
        _must_fail(
            {"cover": {"background_image": "assets/missing.png"}},
            "画像が見つかりません", base_dir=tmp_path,
        )
        _must_fail(
            {"cover": {"background_image": "cover.svg"}},
            "PNGまたはJPEG", base_dir=tmp_path,
        )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        render_cover(
            slide, {"title": "Title", "subtitle": "Subtitle"},
            {"title": "T", "footer": "F", "date": "D", "author": "A"},
            10, loaded,
            add_text=lambda *args, **kwargs: None,
            add_rect=lambda *args, **kwargs: None,
        )
        picture = slide.shapes[0]
        assert picture.name == COVER_BACKGROUND_NAME
        assert picture.width == prs.slide_width
        assert picture.height == prs.slide_height
        assert picture.crop_left > 0
        assert picture.crop_left == picture.crop_right

    print("cover/footer tests passed")


if __name__ == "__main__":
    main()
