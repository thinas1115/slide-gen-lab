"""ハブ図と構成図で共有する描画部品。"""
import math

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from asset_paths import ASSET_DIR, resolve_icon_path
from generate import (ACCENT, BODY_TOP, BODY_BOTTOM, BODY_W, CANVAS, CORAL, GRAY,
                      LIGHT, MARGIN, NAVY, TEXT, WHITE, ContentArea,
                      add_rect, add_text, header, note_line)
from layout_fit import FitError, ensure_within, fit_text_or_raise, select_fit
from textfit import line_height_in, text_width_in

ORANGE = RGBColor(0xE8, 0x7B, 0x1E)   # compute
GREEN = RGBColor(0x3F, 0x86, 0x24)    # storage
PURPLE = RGBColor(0x7D, 0x3F, 0x98)   # ML
LINE = RGBColor(0x6C, 0x73, 0x72)
NODE_LABEL_PAD_X = 0.10  # 左右0.05in。配線以外を隠さない最小限のマスク余白
NODE_LABEL_PAD_Y = 0.04  # 上下0.02in。文字のアンチエイリアス分だけ確保
EDGE_LABEL_PAD_X = 0.10
EDGE_LABEL_PAD_Y = 0.04


def add_arrow(slide, x1, y1, x2, y2, *, color=LINE, width=1.5, both=False,
              dash=None):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    heads = [("a:tailEnd", True), ("a:headEnd", both)]
    for tag, on in heads:
        if on:
            el = ln.makeelement(qn(tag), {"type": "triangle", "w": "med", "len": "med"})
            ln.append(el)
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.insert(0, d)
    return conn


def arrow_label(slide, cx, cy, text, w=1.6, size=9):
    """線をまたぐラベル。白背景マスクは実測テキスト幅に合わせ、wは上限として扱う
    (固定幅だと短い文字列ほど余計な範囲まで線を隠してしまうため)。
    """
    box_h = line_height_in(size, 1.1) + EDGE_LABEL_PAD_Y
    actual_size, _ = fit_text_or_raise(
        "arrow_label", "text", text, w - EDGE_LABEL_PAD_X, box_h, size,
        min_pt=max(7, size - 2), spacing=1.1)
    actual_w = min(w, text_width_in(text, actual_size) + EDGE_LABEL_PAD_X)
    actual_h = line_height_in(actual_size, 1.1) + EDGE_LABEL_PAD_Y
    tb = add_text(slide, cx - actual_w / 2, cy - actual_h / 2, actual_w, actual_h,
                  text, actual_size, color=TEXT, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE)
    tb.fill.solid()
    tb.fill.fore_color.rgb = CANVAS
    return tb


def container(slide, x, y, w, h, label, color=LINE, dash=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.background()
    sp.line.color.rgb = color
    sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    if dash:
        ln = sp.line._get_or_add_ln()
        ln.insert(0, ln.makeelement(qn("a:prstDash"), {"val": dash}))
    label_size, _ = fit_text_or_raise(
        "container", "label", label, w - 0.3, 0.28, 10.5,
        min_pt=8.5, weight="bold", spacing=1.1)
    label_w = min(w - 0.3, text_width_in(label, label_size, "bold") + 0.08)
    add_text(slide, x + 0.12, y + 0.06, label_w, 0.28, label, label_size, bold=True,
             color=color)


def _masked_node_label(slide, cx, y, text, *, max_w, slot_h, size, min_pt,
                       color, bold=False):
    """配線を背後へ通すため、文字の実測外形だけを背景色でマスクする。"""
    weight = "bold" if bold else "regular"
    actual_size, lines = fit_text_or_raise(
        "icon_node", "title" if bold else "sub", text, max_w, slot_h, size,
        min_pt=min_pt, weight=weight, spacing=1.1)
    rendered = "\n".join(lines)
    actual_w = min(max_w, text_width_in(rendered, actual_size, weight)
                   + NODE_LABEL_PAD_X)
    actual_h = min(slot_h, line_height_in(actual_size, 1.1) * len(lines)
                   + NODE_LABEL_PAD_Y)
    tb = add_text(
        slide, cx - actual_w / 2, y + (slot_h - actual_h) / 2,
        actual_w, actual_h, rendered, actual_size, bold=bold, color=color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)
    tb.fill.solid()
    tb.fill.fore_color.rgb = CANVAS
    return tb


# ---- 構成図アイコンノード ----
ICON_R = 0.31        # アイコン半径(0.62角の半分)
EDGE_GAP = 0.06      # 矢印端点とアイコン縁の隙間


def icon_node(slide, cx, cy, img, title, sub=None, size=0.62, *, label_above=False):
    """アイコンと直下のサービス名ラベルを描画する。"""
    p = resolve_icon_path(img)
    if not p.exists():
        raise FileNotFoundError(
            f"アイコン {img} が {ASSET_DIR} にありません。extract_aws_icons.py で"
            f"AWS公式デッキを指定して生成するか、fetch_fluent_icons.py --list で同梱アイコンを確認して"
            f"ください。")
    slide.shapes.add_picture(str(p), Inches(cx - size / 2),
                             Inches(cy - size / 2), Inches(size), Inches(size))
    title_y = cy - size / 2 - 0.59 if label_above else cy + size / 2 + 0.05
    _masked_node_label(slide, cx, title_y, title, max_w=2.1, slot_h=0.28,
                       size=11, min_pt=9.5, color=NAVY, bold=True)
    if sub:
        sub_y = cy - size / 2 - 0.31 if label_above else cy + size / 2 + 0.33
        _masked_node_label(slide, cx, sub_y, sub, max_w=2.1, slot_h=0.26,
                           size=9, min_pt=8, color=GRAY)


def right_of(cx):
    """アイコン右縁の矢印端点X。"""
    return cx + ICON_R + EDGE_GAP


def left_of(cx):
    return cx - ICON_R - EDGE_GAP


# (旧ハンドコード版 s_aws は diagram_layout.py + diagram_specs.py に移行済み)


# ---- ステークホルダー調整図(ハブ型) ----
def fit_hub_layout(area_height, count):
    """放射間隔を圧縮してからアイコンを縮小する。"""
    candidates = [
        ("standard", {"radius_y": 1.72, "icon_size": 0.64}, 4.15),
        ("gap", {"radius_y": 1.54, "icon_size": 0.64}, 3.79),
        ("element", {"radius_y": 1.44, "icon_size": 0.50}, 3.45),
    ]
    if count >= 7:
        candidates = [
            (stage, values, used + 0.10)
            for stage, values, used in candidates
        ]
    return select_fit(
        "hub", area_height, candidates,
        guidance="周辺ノードの文言または件数を減らしてください。",
    )


def _hub_overlap(a, b, gap=0.04):
    return not (a[2] + gap <= b[0] or b[2] + gap <= a[0]
                or a[3] + gap <= b[1] or b[3] + gap <= a[1])


def _hub_label_obstacles(pos, ring, icon_size, cx, cy, hw, hh):
    obstacles = [(cx - hw / 2, cy - hh / 2, cx + hw / 2, cy + hh / 2)]
    for (nx, ny), item in zip(pos, ring):
        obstacles.append((nx - icon_size / 2, ny - icon_size / 2,
                          nx + icon_size / 2, ny + icon_size / 2))
        above = ny > cy + 0.45
        title_y = (ny - icon_size / 2 - 0.59 if above
                   else ny + icon_size / 2 + 0.05)
        title_w = min(2.1, text_width_in(item["name"], 11, "bold")
                      + NODE_LABEL_PAD_X)
        obstacles.append((nx - title_w / 2, title_y,
                          nx + title_w / 2, title_y + 0.28))
        if item.get("sub"):
            sub_y = (ny - icon_size / 2 - 0.31 if above
                     else ny + icon_size / 2 + 0.33)
            sub_w = min(2.1, text_width_in(item["sub"], 9) + NODE_LABEL_PAD_X)
            obstacles.append((nx - sub_w / 2, sub_y,
                              nx + sub_w / 2, sub_y + 0.26))
    return obstacles


def _place_hub_relation_labels(slide, routes, obstacles, area):
    placed = []
    for start_x, start_y, end_x, end_y, label in routes:
        dx, dy = end_x - start_x, end_y - start_y
        distance = math.hypot(dx, dy)
        perp_x, perp_y = -dy / distance, dx / distance
        actual_w = min(2.1, text_width_in(label, 10) + EDGE_LABEL_PAD_X)
        actual_h = line_height_in(10, 1.1) + EDGE_LABEL_PAD_Y
        chosen = None
        for ratio, offset in (
                (0.58, 0.0), (0.64, 0.18), (0.64, -0.18),
                (0.70, 0.24), (0.70, -0.24), (0.52, 0.24), (0.52, -0.24),
                (0.60, 0.42), (0.60, -0.42), (0.55, 0.66), (0.55, -0.66),
                (0.84, 0.0), (0.84, 0.35), (0.84, -0.35)):
            label_x = start_x + dx * ratio + perp_x * offset
            label_y = start_y + dy * ratio + perp_y * offset
            rect = (label_x - actual_w / 2, label_y - actual_h / 2,
                    label_x + actual_w / 2, label_y + actual_h / 2)
            inside = (rect[0] >= MARGIN and rect[2] <= MARGIN + BODY_W
                      and rect[1] >= area.top and rect[3] <= area.bottom)
            if not inside:
                continue
            if any(_hub_overlap(rect, other) for other in obstacles + placed):
                continue
            chosen = (label_x, label_y, rect)
            break
        if chosen is None:
            raise FitError(
                f"hub.ring.label: {label!r} を周辺ノードとハブに重ならず"
                "配置できません。関係ラベルを短くしてください。")
        arrow_label(slide, chosen[0], chosen[1], label, w=2.1, size=10)
        placed.append(chosen[2])


def s_hub(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    if spec.get("note") and area.shifted:
        area = ContentArea(area.top, area.bottom - 0.30, area.shifted)
    y = area.map_y
    cx, cy = 6.67, area.top + area.height * 0.50
    hw, hh = 2.45, 1.05
    ring = spec["ring"]
    if not 3 <= len(ring) <= 8:
        raise FitError(
            "hub: 周辺ノードは3〜8件までです。項目を整理するか図を分割してください。")
    fitted = fit_hub_layout(area.height, len(ring))
    icon_size = fitted.values["icon_size"]
    radius_y = fitted.values["radius_y"]
    radius_x = 4.56
    pos = []
    for index in range(len(ring)):
        angle = -math.pi / 2 + index * 2 * math.pi / len(ring)
        pos.append((cx + radius_x * math.cos(angle),
                    cy + radius_y * math.sin(angle)))

    # Connectors first: lines stay behind the icon and its labels.
    routes = []
    for (sx, sy), item in zip(pos, ring):
        dx, dy = cx - sx, cy - sy
        distance = math.hypot(dx, dy)
        ux, uy = dx / distance, dy / distance
        start_x, start_y = sx + ux * 0.34, sy + uy * 0.34
        ellipse_r = 1 / math.sqrt((ux / (hw / 2)) ** 2 + (uy / (hh / 2)) ** 2)
        end_x, end_y = cx - ux * ellipse_r, cy - uy * ellipse_r
        add_arrow(slide, start_x, start_y, end_x, end_y,
                  both=True, color=LINE, width=1.25)
        routes.append((start_x, start_y, end_x, end_y, item["label"]))

    obstacles = _hub_label_obstacles(pos, ring, icon_size, cx, cy, hw, hh)
    _place_hub_relation_labels(slide, routes, obstacles, area)

    for i, ((nx, ny), item) in enumerate(zip(pos, ring)):
        icon_node(slide, nx, ny, item["icon"], item["name"], item.get("sub"),
                  size=icon_size, label_above=ny > cy + 0.45)

    # 中心ハブ
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - hw / 2), Inches(cy - hh / 2),
                                Inches(hw), Inches(hh))
    sp.fill.solid()
    sp.fill.fore_color.rgb = NAVY
    sp.line.fill.background()
    sp.shadow.inherit = False
    hub_size, hub_lines = fit_text_or_raise(
        "hub", "hub", spec["hub"], hw, 0.62, 13.5,
        min_pt=11, weight="bold", spacing=1.1)
    add_text(slide, cx - hw / 2, cy - 0.31, hw, 0.62,
             spec["hub"], hub_size, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    max_bottom = max(ny + icon_size / 2 + (0.02 if ny > cy + 0.45 else 0.59)
                     for _nx, ny in pos)
    ensure_within(
        "hub", max_bottom - area.top, area.height,
        guidance="周辺ノードのラベルを短くしてください。")
    if spec.get("note"):
        note_line(slide, spec["note"])
