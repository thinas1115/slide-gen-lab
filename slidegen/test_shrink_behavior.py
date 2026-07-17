"""大規模な表・構成図で段階的縮小が実際に発動することを検証する。"""
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

import generate
from content_stress_patterns import (
    LARGE_AWS_DIAGRAM,
    LARGE_DIAGRAM,
    IMAGE_STRESS,
    PROGRAM_ROADMAP_STRESS,
    ROADMAP_STRESS,
    STRESS_PATTERN_DECK,
    TABLE_ROWS,
)
from diagram_layout import ICON_SIZE, Layout
from diagrams2 import PROGRAM_LINE_PT
from generate_from_json import RENDER
from image_slide import fit_image_layout
from timeline_layout import fit_program_roadmap, fit_roadmap, pack_activities
from validate_content import validate


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs.slides.add_slide(prs.slide_layouts[6])


def main():
    errors = validate(deepcopy(STRESS_PATTERN_DECK))
    assert not errors, "\n".join(errors)

    table_spec = STRESS_PATTERN_DECK["slides"][0]
    table_area = generate.header(
        _slide(), table_spec["kicker"], table_spec["title"], table_spec["lead"])
    table_available = (
        table_area.height - generate.TABLE_TOP_GAP - generate.TABLE_BOTTOM_GAP
    )
    table_avail = table_available - generate.TABLE_HEADER_H
    table_widths = generate._auto_table_widths(
        table_spec["columns"], table_spec["rows"])
    table_fit, _row_hs = generate._fit_table(
        TABLE_ROWS, table_widths, table_avail)
    assert table_fit.stage == "font", table_fit
    assert table_fit.values["size"] < 13.5, table_fit
    assert table_fit.values["size"] >= 10.5, table_fit

    table_slide = _slide()
    generate.s_table(table_slide, table_spec, 1)
    table_shape = next(shape for shape in table_slide.shapes if shape.has_table)
    assert table_shape.top + table_shape.height <= Inches(table_area.bottom)

    diagram_spec = STRESS_PATTERN_DECK["slides"][1]
    diagram_area = generate.header(
        _slide(), diagram_spec["kicker"], diagram_spec["title"], diagram_spec["lead"])
    diagram_layout = Layout(deepcopy(LARGE_DIAGRAM), content_area=diagram_area)
    assert len(LARGE_DIAGRAM["nodes"]) == 20
    assert len(LARGE_DIAGRAM["edges"]) >= 20
    assert len(LARGE_DIAGRAM["containers"]) >= 3
    assert any(
        member.startswith("@")
        for container in LARGE_DIAGRAM["containers"]
        for member in container["members"]
    )
    assert diagram_layout.fit_stage == "icon", diagram_layout.fit_stage
    assert diagram_layout.icon_size < ICON_SIZE, diagram_layout.icon_size
    assert diagram_layout.gaps_compressed
    assert max(rect[3] for rect in diagram_layout.cont_rect.values()) \
        <= generate.BODY_BOTTOM
    assert diagram_layout.cont_rect["private"][1] \
        - diagram_layout.cont_rect["public"][3] >= 0.10
    routed = diagram_layout.route_edges(LARGE_DIAGRAM["edges"])
    diagram_layout.validate_edges(LARGE_DIAGRAM["edges"], routed)
    assert any(len(points) >= 4 for points in routed)
    assert any(edge.get("both") for edge in LARGE_DIAGRAM["edges"])
    assert any(edge.get("dash") for edge in LARGE_DIAGRAM["edges"])
    assert any(len(edge.get("via", [])) > 0 for edge in LARGE_DIAGRAM["edges"])

    aws_spec = STRESS_PATTERN_DECK["slides"][2]
    aws_area = generate.header(
        _slide(), aws_spec["kicker"], aws_spec["title"], aws_spec["lead"])
    aws_layout = Layout(deepcopy(LARGE_AWS_DIAGRAM), content_area=aws_area)
    assert len(LARGE_AWS_DIAGRAM["nodes"]) >= 16
    assert len(LARGE_AWS_DIAGRAM["containers"]) >= 3
    assert all(
        not node["icon"].startswith("icons/fluent/")
        for node in LARGE_AWS_DIAGRAM["nodes"].values()
    )
    assert aws_layout.fit_stage == "icon", aws_layout.fit_stage
    assert aws_layout.icon_size < ICON_SIZE, aws_layout.icon_size
    assert aws_layout.gaps_compressed
    assert max(rect[3] for rect in aws_layout.cont_rect.values()) \
        <= generate.BODY_BOTTOM
    aws_routed = aws_layout.route_edges(LARGE_AWS_DIAGRAM["edges"])
    aws_layout.validate_edges(LARGE_AWS_DIAGRAM["edges"], aws_routed)
    worker_routes = {
        edge["to"]: points
        for edge, points in zip(LARGE_AWS_DIAGRAM["edges"], aws_routed)
        if edge["from"] == "worker_b"
    }
    ddb_start, ddb_next = worker_routes["ddb"][:2]
    assert ddb_next[0] > ddb_start[0]
    assert abs(ddb_next[1] - ddb_start[1]) <= 0.01
    rds_start, rds_next = worker_routes["rds_b"][:2]
    assert abs(rds_next[0] - rds_start[0]) <= 0.01
    assert rds_next[1] > rds_start[1]
    assert len(worker_routes["rds_b"]) == 2
    assert any(len(points) >= 4 for points in aws_routed)
    assert any(edge.get("both") for edge in LARGE_AWS_DIAGRAM["edges"])
    assert any(edge.get("dash") for edge in LARGE_AWS_DIAGRAM["edges"])
    assert any(len(edge.get("via", [])) > 0 for edge in LARGE_AWS_DIAGRAM["edges"])

    bad_edges = [
        {"from": "a", "to": "b"},
        {"from": "c", "to": "d"},
    ]
    bad_routes = [
        [(0.0, 0.0), (2.0, 0.0)],
        [(1.0, -1.0), (1.0, 0.0)],
    ]
    try:
        aws_layout._validate_edge_contacts(bad_edges, bad_routes)
    except ValueError as exc:
        assert "接続して見える" in str(exc)
    else:
        raise AssertionError("別エッジのT字接触が検出されませんでした")

    fanout_edges = [
        {"from": "a", "to": "b"},
        {"from": "a", "to": "c"},
    ]
    aws_layout._validate_edge_contacts(fanout_edges, [
        [(0.0, 0.0), (1.0, 0.0), (1.0, 1.0)],
        [(0.0, 0.0), (1.0, 0.0), (1.0, -1.0)],
    ])
    try:
        aws_layout._validate_edge_contacts(fanout_edges, [
            [(0.0, 0.0), (0.0, 1.0), (2.0, 1.0)],
            [(0.0, 0.0), (1.0, 0.0), (1.0, 2.0)],
        ])
    except ValueError as exc:
        assert "接続して見える" in str(exc)
    else:
        raise AssertionError("同一始点から分岐した後の再交差が検出されませんでした")

    try:
        aws_layout._validate_segment_lengths(
            [{"from": "user", "to": "r53"}],
            [[(0.0, 0.0), (1.0, 1.0)]],
        )
    except ValueError as exc:
        assert "直角ではありません" in str(exc)
    else:
        raise AssertionError("斜めの経路区間が検出されませんでした")

    try:
        aws_layout._validate_segment_lengths(
            [{"from": "user", "to": "r53", "exit": "bottom", "enter": "top"}],
            [[(0.0, 0.0), (1.0, 0.0), (1.0, 1.0)]],
        )
    except ValueError as exc:
        assert "垂直外向きではありません" in str(exc)
    else:
        raise AssertionError("接続辺に対して垂直でない開始区間が検出されませんでした")

    roadmap_area = generate.header(
        _slide(), ROADMAP_STRESS["kicker"], ROADMAP_STRESS["title"],
        ROADMAP_STRESS["lead"])
    roadmap_fit = fit_roadmap(
        roadmap_area.height, len(ROADMAP_STRESS["phases"]), has_note=True)
    assert roadmap_fit.stage == "element", roadmap_fit
    assert roadmap_fit.values["row_h"] < 0.70, roadmap_fit

    program_area = generate.header(
        _slide(), PROGRAM_ROADMAP_STRESS["kicker"],
        PROGRAM_ROADMAP_STRESS["title"], PROGRAM_ROADMAP_STRESS["lead"])
    lane_counts = [
        pack_activities(track["activities"], PROGRAM_ROADMAP_STRESS["periods"])[1]
        for track in PROGRAM_ROADMAP_STRESS["tracks"]
    ]
    assert lane_counts == [3, 3, 3, 3, 3]
    program_fit = fit_program_roadmap(program_area.height, lane_counts)
    assert program_fit.stage == "element", program_fit
    assert program_fit.values["lane_pitch"] < 0.30, program_fit

    program_slide = _slide()
    generate.render_slide(
        RENDER["program_roadmap"], program_slide, PROGRAM_ROADMAP_STRESS, 1)
    track_names = {track["name"] for track in PROGRAM_ROADMAP_STRESS["tracks"]}
    track_shapes = [
        shape for shape in program_slide.shapes
        if shape.has_text_frame and shape.text_frame.text in track_names
    ]
    assert len(track_shapes) == len(track_names)
    assert all(shape.text_frame.word_wrap is False for shape in track_shapes)
    assert max(shape.width.inches for shape in track_shapes) - min(
        shape.width.inches for shape in track_shapes) < 0.01
    activity_lines = [
        shape for shape in program_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.LINE
        and shape.line.width is not None
        and abs(shape.line.width.pt - PROGRAM_LINE_PT) < 0.01
    ]
    assert len(activity_lines) == 15

    short_spec = deepcopy(PROGRAM_ROADMAP_STRESS)
    short_spec["tracks"] = [{
        "name": "短期作業",
        "activities": [
            {"label": "半月作業", "start": 4.25, "end": 4.75},
            {"label": "後続作業", "start": 5.0, "end": 6.0},
        ],
    }]
    short_slide = _slide()
    generate.render_slide(
        RENDER["program_roadmap"], short_slide, short_spec, 1)
    short_label = next(
        shape for shape in short_slide.shapes
        if shape.has_text_frame and shape.text_frame.text == "半月作業")
    short_lines = sorted(
        (
            shape for shape in short_slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.LINE
            and shape.line.width is not None
            and abs(shape.line.width.pt - PROGRAM_LINE_PT) < 0.01
        ),
        key=lambda shape: shape.left,
    )
    assert len(short_lines) == 2
    short_line = short_lines[0]
    assert abs(
        (short_label.left + short_label.width / 2)
        - (short_line.left + short_line.width / 2)
    ) < Inches(0.01)

    image_area = generate.header(
        _slide(), IMAGE_STRESS["kicker"], IMAGE_STRESS["title"],
        IMAGE_STRESS["lead"])
    image_fit = fit_image_layout(min(image_area.height, 3.15))
    assert image_fit.stage == "element", image_fit
    assert image_fit.values["min_image_h"] < 3.20, image_fit

    generate.DECK = STRESS_PATTERN_DECK
    for idx, spec in enumerate(STRESS_PATTERN_DECK["slides"], 1):
        generate.render_slide(RENDER[spec["type"]], _slide(), spec, idx)

    print(
        "shrink behavior tests passed "
        f"(table={table_fit.values['size']:.1f}pt, "
        f"fluent={diagram_layout.icon_size:.2f}in, "
        f"aws={aws_layout.icon_size:.2f}in / {ICON_SIZE:.2f}in, "
        f"roadmap={roadmap_fit.values['row_h']:.2f}in, "
        f"program={program_fit.values['lane_pitch']:.2f}in, "
        f"image={image_fit.values['min_image_h']:.2f}in)"
    )


if __name__ == "__main__":
    main()
