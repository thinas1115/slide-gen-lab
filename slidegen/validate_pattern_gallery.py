"""パターン検証ギャラリーの結果をスライド別一覧にする。"""
import csv
import sys
from collections import defaultdict
from pathlib import Path

from PIL import Image
from pptx import Presentation

from check_layout import check
from content_patterns import PATTERN_DECK


def png_status(png_dir, idx):
    p = Path(png_dir) / f"slide_{idx:02d}.png"
    if not p.exists():
        return "未生成"
    with Image.open(p) as img:
        return f"OK ({img.width}x{img.height})"


def write_reports(rows, md_path, csv_path):
    md = [
        "# パターン検証結果",
        "",
        "| # | パターン | タイトル | レイアウト検査 | PNG |",
        "|---:|---|---|---|---|",
    ]
    for r in rows:
        md.append("| {idx} | {pattern} | {title} | {layout} | {png} |".format(**r))
    md_path = Path(md_path)
    csv_path = Path(csv_path)
    md_path.parent.mkdir(parents=True, exist_ok=True)
    csv_path.parent.mkdir(parents=True, exist_ok=True)
    md_path.write_text("\n".join(md) + "\n", encoding="utf-8")

    with csv_path.open("w", newline="", encoding="utf-8-sig") as f:
        writer = csv.DictWriter(f, fieldnames=["idx", "pattern", "title", "layout", "png"])
        writer.writeheader()
        writer.writerows(rows)
    print(f"出力: {md_path}")
    print(f"出力: {csv_path}")


def main(pptx_path, png_dir, md_path, csv_path):
    prs = Presentation(pptx_path)
    specs = PATTERN_DECK["slides"]
    if len(prs.slides) != len(specs):
        raise SystemExit(
            f"スライド数不一致: PPTX={len(prs.slides)}、仕様={len(specs)}")

    by_slide = defaultdict(list)
    for si, kind, a, b in check(pptx_path):
        by_slide[si].append(f"{kind}: {a} x {b}")

    rows = []
    for idx, spec in enumerate(specs, 1):
        findings = by_slide.get(idx, [])
        rows.append({
            "idx": idx,
            "pattern": spec.get("pattern", spec["type"]),
            "title": spec.get("title", "").replace("|", "｜").replace("\n", "<br>"),
            "layout": "OK" if not findings else "<br>".join(findings).replace("|", "｜"),
            "png": png_status(png_dir, idx),
        })
    write_reports(rows, md_path, csv_path)


if __name__ == "__main__":
    project_root = Path(__file__).resolve().parent.parent
    out_dir = project_root / "out"
    main(
        sys.argv[1] if len(sys.argv) > 1 else out_dir / "pattern_gallery.pptx",
        sys.argv[2] if len(sys.argv) > 2 else out_dir / "png_pattern_gallery",
        sys.argv[3] if len(sys.argv) > 3 else out_dir / "pattern_gallery_validation.md",
        sys.argv[4] if len(sys.argv) > 4 else out_dir / "pattern_gallery_validation.csv",
    )
