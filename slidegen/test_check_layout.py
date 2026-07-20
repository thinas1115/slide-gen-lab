"""check_layoutが壊れたPPTXを誤って合格させないための回帰テスト。"""
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from check_layout import check


def save(prs, path):
    prs.save(path)
    return check(path)


with TemporaryDirectory() as td:
    out = Path(td)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_table(1, 2, Inches(1), Inches(7.1), Inches(8), Inches(1))
    findings = save(prs, out / "table_oob.pptx")
    assert any(kind == "OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = ChartData()
    data.categories = ["A", "B"]
    data.add_series("値", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(12.8), Inches(1),
        Inches(1), Inches(2), data)
    findings = save(prs, out / "chart_oob.pptx")
    assert any(kind == "OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(9.8), Inches(1), Inches(0.5), Inches(0.5))
    findings = save(prs, out / "custom_size_oob.pptx")
    assert any(kind == "OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(
        1, 1, Inches(1), Inches(1), Inches(1.2), Inches(0.25)).table
    cell = table.cell(0, 0)
    cell.text = "セルから確実にはみ出す長い文章です"
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(20)
    findings = save(prs, out / "cell_text_oob.pptx")
    assert any(kind == "CELL-OOB" for _, kind, _, _ in findings), findings

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(5), Inches(1.5)).table
    merged = table.cell(0, 0)
    merged.merge(table.cell(0, 1))
    merged.text = "結合セルの幅全体を使う見出し"
    findings = save(prs, out / "merged_table_ok.pptx")
    assert not any(kind == "CELL-OOB" for _, kind, _, _ in findings), findings

print("check_layout broken-PPTX regression: ALL OK")
