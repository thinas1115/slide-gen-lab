"""基本デッキに表現力検証スライドを加えた拡張サンプルを生成する。"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import generate
from content import DECK
from content_ext import EXTRA_SLIDES
from diagrams import s_hub, s_org
from diagrams2 import s_matrix, s_process, s_roadmap
from diagram_layout import render_diagram


def s_diagram(slide, spec, page):
    """content.json と同じインライン仕様による図解スライド。"""
    generate.header(slide, spec["kicker"], spec["title"])
    render_diagram(slide, spec["diagram"], note=spec.get("note"))


RENDER_EXTENDED = dict(generate.RENDER,
                       hub=s_hub, org=s_org,
                       process=s_process, roadmap=s_roadmap, matrix=s_matrix,
                       diagram=s_diagram)


def main(out_path, cover_footer_config=None):
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        generate.configure_cover_footer(cover_footer_config)
    except ValueError as e:
        raise SystemExit(f"NG: 表紙・フッター設定: {e}") from e
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    slides = DECK["slides"] + EXTRA_SLIDES
    generate.DECK = {**DECK, "slides": slides}
    for idx, spec in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        RENDER_EXTENDED[spec["type"]](slide, spec, idx)
        if spec["type"] != "title":
            generate.footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(slides)} slides)")


if __name__ == "__main__":
    default_out = Path(__file__).resolve().parent.parent / "out" / "sample_extended.pptx"
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out_path", nargs="?", default=default_out)
    parser.add_argument("--cover-footer-config", metavar="PATH",
                        help="表紙・フッター設定JSON")
    args = parser.parse_args()
    main(args.out_path, args.cover_footer_config)
