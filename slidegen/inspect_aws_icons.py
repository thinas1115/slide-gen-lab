"""AWSアイコンデッキの構造調査: アイコンが画像として抽出できるか確認する。"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

if len(sys.argv) != 2:
    raise SystemExit(
        "使い方: python slidegen/inspect_aws_icons.py <AWS公式アイコンデッキ.pptx>")
src = Path(sys.argv[1]).expanduser()
if not src.is_file():
    raise SystemExit(f"入力PPTXが見つかりません: {src}")
out = []

prs = Presentation(src)
out.append(f"slides: {len(prs.slides)}")


def walk(shapes, depth=0):
    for sh in shapes:
        yield sh, depth
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk(sh.shapes, depth + 1)


# 最初の10枚でshape種別の分布とnameの様子を見る
from collections import Counter
kinds = Counter()
samples = []
for i, slide in enumerate(prs.slides):
    if i >= 10:
        break
    for sh, d in walk(slide.shapes):
        kinds[str(sh.shape_type)] += 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and len(samples) < 15:
            samples.append(f"  s{i + 1} pic name={sh.name!r} size={sh.width}x{sh.height}")

out.append("shape types (first 10 slides): " + str(dict(kinds)))
out += samples

# 対象サービス名がnameに含まれるPictureを全スライドから検索
targets = ["Bedrock", "Simple Storage", "S3", "Fargate", "Elastic Container",
           "CloudWatch", "Elastic Load", "Application Load", "Virtual-private",
           "VPC", "Users", "User"]
hits = {}
for i, slide in enumerate(prs.slides):
    for sh, d in walk(slide.shapes):
        if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        nm = sh.name
        for t in targets:
            if t.lower() in nm.lower() and t not in hits:
                hits[t] = f"s{i + 1}: {nm!r}"
out.append("target hits:")
out += [f"  {k}: {v}" for k, v in hits.items()]

Path(__file__).with_name("aws_icons_report.txt").write_text(
    "\n".join(out), encoding="utf-8")
print("done")
