"""全rendererが過密入力を明示停止することを検証する。"""
from pptx import Presentation
from pptx.util import Inches

import generate
from diagrams import s_hub
from org_layout import s_org
from diagrams2 import s_matrix, s_process, s_program_roadmap, s_roadmap
from layout_fit import FitError


LONG = "提出品質を維持できないほど長い説明文です。" * 10


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _must_fail(renderer, spec, expected):
    try:
        renderer(_slide(), spec, 1)
    except FitError as e:
        assert expected in str(e), str(e)
    else:
        raise AssertionError(f"{spec['type']}が過密入力を拒否しませんでした")


def _base(type_):
    return {"type": type_, "kicker": "TEST", "title": "収容検証"}


def main():
    spec = _base("bullets")
    spec["bullets"] = [[LONG, None] for _ in range(6)]
    _must_fail(generate.s_bullets, spec, "不足")
    try:
        generate.render_slide(generate.s_bullets, _slide(), spec, 1)
    except SystemExit as e:
        assert "slides[0] (type=bullets)" in str(e), str(e)
        assert "箇条書きを減らす" in str(e), str(e)
    else:
        raise AssertionError("rendererエラーにスライド位置が付与されませんでした")

    spec = _base("cards")
    spec.update(style="editorial", cards=[[f"項目{i}", LONG] for i in range(7)])
    _must_fail(generate.s_cards, spec, "カード本文")

    spec = _base("table")
    spec.update(columns=["項目", "説明"],
                rows=[[f"行{i}", LONG] for i in range(8)])
    _must_fail(generate.s_table, spec, "表の行を減らす")

    panel = {"heading": "比較", "bullets": [LONG for _ in range(6)]}
    spec = _base("twocol")
    spec.update(left=panel, right=panel)
    _must_fail(generate.s_twocol, spec, "左右の箇条書き")

    spec = _base("chart")
    spec["chart"] = {
        "categories": [str(i) for i in range(13)],
        "series": [["実績", [1] * 13]],
    }
    _must_fail(generate.s_chart, spec, "カテゴリ1〜12件")

    spec = _base("process")
    spec.update(steps=[{"name": "工程", "desc": "説明", "actor": "担当"}
                       for _ in range(7)], emph=[])
    _must_fail(s_process, spec, "工程は3〜6件")

    optional_actor = _base("process")
    optional_actor.update(steps=[
        {"name": "工程A", "desc": "担当表示なし"},
        {"name": "工程B", "desc": "担当表示あり", "actor": "担当区分"},
        {"name": "工程C", "desc": "担当表示なし"},
    ])
    s_process(_slide(), optional_actor, 1)

    spec = _base("roadmap")
    spec.update(months=["1月", "2月", "3月", "4月"], milestones=[],
                phases=[{"name": "Phase", "goal": "目標", "bar": "実行",
                         "start": 0, "end": 1} for _ in range(7)])
    _must_fail(s_roadmap, spec, "フェーズは1〜6件")

    spec = _base("program_roadmap")
    spec.update(periods=["1月", "2月", "3月", "4月"],
                tracks=[{"name": "テーマ", "activities": [
                    {"label": f"作業{j}", "start": 0, "end": 4}
                    for j in range(4)
                ]} for _ in range(6)])
    _must_fail(s_program_roadmap, spec, "最小設定")

    spec = _base("matrix")
    spec.update(x_axis="X", y_axis="Y", target_label="対象",
                points=[{"name": str(i), "x": 0.5, "y": 0.5}
                        for i in range(9)])
    _must_fail(s_matrix, spec, "点は1〜8件")

    spec = _base("hub")
    spec.update(hub="中心", ring=[{"name": "部門", "label": "連携",
                                   "icon": "icons/fluent/team.png"}
                                  for _ in range(9)])
    _must_fail(s_hub, spec, "周辺ノードは3〜8件")

    spec = _base("org")
    spec["org"] = {
        "nodes": {
            f"level{i}": {"name": f"第{i + 1}階層", "sub": "担当範囲",
                          "members": ["担当A", "担当B"]}
            for i in range(6)
        },
        "levels": [[f"level{i}"] for i in range(6)],
        "edges": [{"from": f"level{i}", "to": f"level{i + 1}"}
                  for i in range(5)],
    }
    _must_fail(s_org, spec, "最小設定")

    print("renderer fit tests passed")


if __name__ == "__main__":
    main()
