"""汎用化したrendererの構造バリエーションと互換入力を検証する。"""
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

import generate
from content_patterns import PATTERN_DECK
from diagrams import fit_hub_layout, s_hub
from diagrams2 import (_fit_process_flow, _matrix_label_positions,
                       fit_matrix_labels, s_matrix, s_process)
from layout_fit import FitError
from validate_content import validate


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _base(type_):
    return {"type": type_, "kicker": "TEST", "title": "汎用化検証"}


def _must_fail(callable_, expected):
    try:
        callable_()
    except FitError as exc:
        assert expected in str(exc), str(exc)
    else:
        raise AssertionError("最小設定での明示停止が発生しませんでした")


def _text_shape(slide, text):
    return next(
        shape for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text == text)


def _assert_process_natural_top(slide, flow, area_top):
    first_row_ids = [level[0] for level in flow["levels"]]
    first_row_tops = [
        _text_shape(slide, flow["nodes"][node_id]["name"]).top / Inches(1)
        for node_id in first_row_ids
    ]
    expected_title_top = area_top + 0.34 + 0.13
    assert max(first_row_tops) - min(first_row_tops) < 0.01
    assert abs(first_row_tops[0] - expected_title_top) < 0.01

    connector_bottoms = [
        (shape.top + shape.height) / Inches(1)
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.LINE
    ]
    assert connector_bottoms
    assert max(connector_bottoms) < generate.BODY_BOTTOM - 0.80


def main():
    errors = validate(deepcopy(PATTERN_DECK), allow_sample_content=True)
    assert not errors, "\n".join(errors)

    columns = ["区分", "短い値", "詳細説明"]
    rows = [["A", "可", "利用部門と運用条件を文章で説明する"]]
    widths = generate._auto_table_widths(columns, rows)
    assert abs(sum(widths) - generate.BODY_W) < 0.001
    assert widths[2] > widths[0]
    table = dict(_base("table"), columns=columns, rows=rows)
    generate.s_table(_slide(), table, 1)

    cards = dict(
        _base("cards"), style="metrics",
        cards=[
            {"heading": f"指標{i + 1}", "value": f"{(i + 1) * 10}%",
             "body": "判断に必要な補足説明", "emphasis": i == 0}
            for i in range(5)
        ],
    )
    card_slide = _slide()
    generate.s_cards(card_slide, cards, 1)
    assert any(shape.has_text_frame and shape.text_frame.text == "10%"
               for shape in card_slide.shapes)

    comparison = dict(
        _base("twocol"),
        left={"label": "現状", "heading": "現在", "bullets": ["課題を確認"]},
        right={"label": "目標", "heading": "将来", "bullets": ["改善を実施"]},
    )
    comparison_slide = _slide()
    generate.s_twocol(comparison_slide, comparison, 1)
    texts = {shape.text_frame.text for shape in comparison_slide.shapes
             if shape.has_text_frame}
    assert {"現状", "目標"} <= texts

    for kind in ("bar", "column", "line", "stacked_bar", "stacked_column"):
        chart = dict(
            _base("chart"),
            chart={
                "kind": kind,
                "categories": ["A", "B", "C"],
                "series": [["実績", [10, 20, 30]], ["計画", [12, 18, 32]]],
            },
        )
        chart_slide = _slide()
        generate.s_chart(chart_slide, chart, 1)
        assert any(shape.has_chart for shape in chart_slide.shapes)
    assert generate._fit_chart_layout(
        generate.ContentArea(), 12, 2, "bar").stage == "element"
    _must_fail(
        lambda: generate._fit_chart_layout(
            generate.ContentArea(2.0, 3.3, True), 12, 4, "bar"),
        "最小設定",
    )

    flow = {
        "nodes": {
            "start": {"name": "受付"},
            "review": {"name": "審査", "style": "decision"},
            "ok": {"name": "承認", "style": "accent"},
            "fix": {"name": "修正"},
            "end": {"name": "実施"},
        },
        "levels": [["start"], ["review"], ["ok", "fix"], ["end"]],
        "edges": [
            {"from": "start", "to": "review"},
            {"from": "review", "to": "ok", "label": "可"},
            {"from": "review", "to": "fix", "label": "修正"},
            {"from": "ok", "to": "end"},
            {"from": "fix", "to": "review", "kind": "feedback"},
        ],
    }
    process_slide = _slide()
    s_process(process_slide, dict(_base("process"), flow=flow), 1)
    _assert_process_natural_top(process_slide, flow, generate.BODY_TOP)

    lead = "審査結果と次の対応を先に説明します。"
    lead_probe = _slide()
    lead_area = generate.header(lead_probe, "TEST", "汎用化検証", lead)
    lead_process_slide = _slide()
    s_process(
        lead_process_slide,
        dict(_base("process"), flow=flow, lead=lead),
        1,
    )
    _assert_process_natural_top(lead_process_slide, flow, lead_area.top)
    dense_levels = [[f"n{column}_{row}" for row in range(3)]
                    for column in range(3)]
    assert _fit_process_flow(4.0, dense_levels).stage == "gap"
    assert _fit_process_flow(3.5, dense_levels).stage == "element"
    _must_fail(lambda: _fit_process_flow(3.0, dense_levels), "最小設定")

    matrix = next(slide for slide in PATTERN_DECK["slides"]
                  if slide["type"] == "matrix")
    assert all("lx" not in point and "ly" not in point
               for point in matrix["points"])
    stage, _width, _xy, labels = fit_matrix_labels(
        matrix["points"], 1.68, 2.22, 10.5, 4.02,
        barriers=(1.68 + 10.5 / 2, 2.22 + 4.02 / 2))
    assert stage in {"standard", "gap", "element"}
    for index, rect in labels.items():
        for other_index, other in labels.items():
            if index < other_index:
                assert not (rect[0] < other[2] and other[0] < rect[2]
                            and rect[1] < other[3] and other[1] < rect[3])
        assert not rect[0] < 1.68 + 10.5 / 2 < rect[2]
        assert not rect[1] < 2.22 + 4.02 / 2 < rect[3]
    legacy_points = [{"name": "互換", "x": 0.4, "y": 0.5,
                      "lx": 0.2, "ly": -0.36}]
    _xy, legacy_labels = _matrix_label_positions(
        legacy_points, 1.0, 2.0, 8.0, 4.0, label_w=2.0)
    legacy_point_x = 1.0 + 0.4 * 8.0
    assert abs(legacy_labels[0][0] - (legacy_point_x - 1.0 + 0.2)) < 0.001
    s_matrix(_slide(), matrix, 1)

    for count in (3, 5, 8):
        hub = dict(
            _base("hub"), hub="中心",
            ring=[{"name": f"部門{i + 1}", "label": "連携",
                   "icon": "icons/fluent/team.png"} for i in range(count)],
        )
        s_hub(_slide(), hub, 1)
    assert fit_hub_layout(4.0, 5).stage == "gap"
    assert fit_hub_layout(3.6, 8).stage == "element"
    _must_fail(lambda: fit_hub_layout(3.2, 8), "最小設定")

    print("generalized renderer tests passed")


if __name__ == "__main__":
    main()
