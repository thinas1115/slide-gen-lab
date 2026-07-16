"""大判画像を本文領域の主役として配置するrenderer。"""
from PIL import Image
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from asset_paths import resolve_image_path
from generate import BODY_W, GRAY, TEXT, add_text, header
from layout_fit import select_fit, stepped
from textfit import line_height_in, wrap_text

FRAME_X = 0.72
FRAME_W = BODY_W - 0.34
MIN_IMAGE_H = 2.40
PICTURE_NAME = "ContentImage"


def _measure_text(text, width, size, spacing):
    if not text:
        return [], 0.0
    lines = wrap_text(text, width, size)
    return lines, len(lines) * line_height_in(size, spacing) + 0.02


def fit_image_layout(available, caption=None, source=None):
    """標準余白、余白圧縮、文字・画像縮小、停止の順で収容する。"""
    candidates = []

    def add(stage, values):
        caption_lines, caption_h = _measure_text(
            caption, FRAME_W, values["caption_pt"], 1.12)
        source_lines, source_h = _measure_text(
            source, FRAME_W, values["source_pt"], 1.08)
        text_h = 0.0
        if caption_lines:
            text_h += values["caption_gap"] + caption_h
        if source_lines:
            text_h += values["source_gap"] + source_h
        values = dict(values, caption_lines=caption_lines,
                      source_lines=source_lines, caption_h=caption_h,
                      source_h=source_h, text_h=text_h)
        used = (values["top_gap"] + values["min_image_h"] + text_h
                + values["bottom_gap"])
        candidates.append((stage, values, used))

    add("standard", {
        "top_gap": 0.16, "bottom_gap": 0.05,
        "caption_gap": 0.10, "source_gap": 0.06,
        "caption_pt": 10.5, "source_pt": 8.0, "min_image_h": 3.20,
    })
    add("gap", {
        "top_gap": 0.08, "bottom_gap": 0.03,
        "caption_gap": 0.06, "source_gap": 0.04,
        "caption_pt": 10.5, "source_pt": 8.0, "min_image_h": 2.90,
    })
    for min_image_h in stepped(2.80, MIN_IMAGE_H, 0.10):
        ratio = (min_image_h - MIN_IMAGE_H) / (2.80 - MIN_IMAGE_H)
        add("element", {
            "top_gap": 0.06, "bottom_gap": 0.02,
            "caption_gap": 0.05, "source_gap": 0.03,
            "caption_pt": 8.5 + ratio * 1.5,
            "source_pt": 7.0 + ratio * 0.5,
            "min_image_h": min_image_h,
        })
    return select_fit(
        "image", available, candidates,
        guidance=("キャプション・出典を短くするか、leadを外すか、"
                  "画像だけのスライドへ分割してください。"),
    )


def _add_picture(slide, path, x, y, w, h, fit, alt=None):
    with Image.open(path) as image:
        image_w, image_h = image.size
    image_ratio = image_w / image_h
    frame_ratio = w / h

    if fit == "cover":
        picture = slide.shapes.add_picture(
            str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        if image_ratio > frame_ratio:
            crop = (1 - frame_ratio / image_ratio) / 2
            picture.crop_left = picture.crop_right = crop
        elif image_ratio < frame_ratio:
            crop = (1 - image_ratio / frame_ratio) / 2
            picture.crop_top = picture.crop_bottom = crop
    else:
        if image_ratio >= frame_ratio:
            picture_w, picture_h = w, w / image_ratio
        else:
            picture_w, picture_h = h * image_ratio, h
        picture = slide.shapes.add_picture(
            str(path), Inches(x + (w - picture_w) / 2),
            Inches(y + (h - picture_h) / 2),
            Inches(picture_w), Inches(picture_h))

    picture.name = PICTURE_NAME
    if alt:
        picture._element.nvPicPr.cNvPr.set("descr", alt)
    return picture


def s_image(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    path = resolve_image_path(spec["image"])
    if not path.is_file():
        raise FileNotFoundError(f"画像 {spec['image']!r} がassets内にありません")

    fitted = fit_image_layout(
        area.height, spec.get("caption"), spec.get("source"))
    values = fitted.values
    top = area.top + values["top_gap"]
    image_h = (area.height - values["top_gap"] - values["bottom_gap"]
               - values["text_h"])
    _add_picture(
        slide, path, FRAME_X, top, FRAME_W, image_h,
        spec.get("fit", "contain"), spec.get("alt"))

    cursor = top + image_h
    if values["caption_lines"]:
        cursor += values["caption_gap"]
        add_text(
            slide, FRAME_X, cursor, FRAME_W, values["caption_h"],
            "\n".join(values["caption_lines"]), values["caption_pt"],
            color=TEXT, spacing=1.12)
        cursor += values["caption_h"]
    if values["source_lines"]:
        cursor += values["source_gap"]
        add_text(
            slide, FRAME_X, cursor, FRAME_W, values["source_h"],
            "\n".join(values["source_lines"]), values["source_pt"],
            color=GRAY, align=PP_ALIGN.RIGHT, spacing=1.08)
