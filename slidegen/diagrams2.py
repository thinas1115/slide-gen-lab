"""図解系スライド第2弾: プロセスタイムライン・ロードマップ・2軸マップ。"""
import re

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from generate import (ACCENT, CORAL, GRAY, LIGHT, NAVY, RULE, TEXT, WHITE,
                      ZEBRA, ContentArea, add_rect, add_text, header, note_line)
from diagrams import add_arrow, arrow_label
from diagrams3 import route
from layout_fit import FitError, ensure_within, fit_text_or_raise, select_fit
from textfit import text_width_in
from timeline_layout import (centered_label_box, fit_program_roadmap,
                             fit_roadmap, pack_activities, resolve_marker,
                             resolve_span)

PROGRAM_LINE_PT = 1.4


def _rects_overlap(a, b, gap=0.04):
    return not (a[2] + gap <= b[0] or b[2] + gap <= a[0]
                or a[3] + gap <= b[1] or b[3] + gap <= a[1])


def _matrix_label_positions(
        points, ox, top, aw, ah, *, label_w=2.0, gap=0.04, barriers=()):
    """点ラベルを周囲8方向へ自動配置し、衝突しない矩形を返す。"""
    point_xy = [(ox + point["x"] * aw, top + (1 - point["y"]) * ah)
                for point in points]
    point_boxes = [(x - 0.17, y - 0.17, x + 0.17, y + 0.17)
                   for x, y in point_xy]
    placed = {}
    order = sorted(
        range(len(points)),
        key=lambda i: (
            -sum(abs(point_xy[i][0] - x) < 2.2 and abs(point_xy[i][1] - y) < 0.7
                 for x, y in point_xy),
            -int(bool(points[i].get("emph"))),
            i,
        ),
    )
    label_h = 0.31
    offsets = [
        (-label_w / 2, -0.49), (-label_w / 2, 0.18),
        (0.22, -label_h / 2), (-label_w - 0.22, -label_h / 2),
        (0.18, -0.43), (-label_w - 0.18, -0.43),
        (0.18, 0.14), (-label_w - 0.18, 0.14),
    ]
    for index in order:
        px, py = point_xy[index]
        manual = (points[index].get("lx"), points[index].get("ly"))
        candidates = offsets
        if manual[0] is not None or manual[1] is not None:
            candidates = [(-label_w / 2 + (manual[0] or 0.0),
                           manual[1] if manual[1] is not None else -0.36)]
        for dx, dy in candidates:
            rect = (px + dx, py + dy, px + dx + label_w, py + dy + label_h)
            inside = (rect[0] >= ox + 0.04 and rect[2] <= ox + aw - 0.04
                      and rect[1] >= top + 0.04 and rect[3] <= top + ah - 0.04)
            if not inside:
                continue
            if any(rect[0] < barrier < rect[2] for barrier in barriers[:1]):
                continue
            if any(rect[1] < barrier < rect[3] for barrier in barriers[1:2]):
                continue
            if any(_rects_overlap(rect, other, gap)
                   for other_index, other in enumerate(point_boxes)
                   if other_index != index):
                continue
            if any(_rects_overlap(rect, other, gap) for other in placed.values()):
                continue
            placed[index] = rect
            break
        else:
            raise FitError(
                f"matrix.points[{index}]: {points[index]['name']!r} のラベルを"
                "他の点と重ならずに配置できません。点を減らすか座標の粒度を"
                "見直してください。")
    return point_xy, placed


def fit_matrix_labels(points, ox, top, aw, ah, *, barriers=()):
    """ラベル間隔を圧縮し、その後ラベル幅を縮めて配置を試す。"""
    last_error = None
    for stage, width, gap in (
            ("standard", 2.0, 0.05),
            ("gap", 2.0, 0.015),
            ("element", 1.65, 0.015)):
        try:
            point_xy, labels = _matrix_label_positions(
                points, ox, top, aw, ah, label_w=width, gap=gap,
                barriers=barriers)
            return stage, width, point_xy, labels
        except FitError as exc:
            last_error = exc
    raise FitError(
        f"matrix: 最小ラベル幅でも衝突を解消できません。{last_error}")


def _grid_line(slide, x1, y1, x2, y2):
    """表の罫線を本文より先に描き、後続のラベルマスクを有効にする。"""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RULE
    line.line.width = Pt(0.6)
    line.shadow.inherit = False
    return line


def _activity_line(slide, x1, y, x2, color):
    """工程表の作業期間を描く。矩形ではなく固定pt幅で見え方を揃える。"""
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(PROGRAM_LINE_PT)
    line.shadow.inherit = False
    return line


def _fit_single_line(renderer, field, text, width, max_pt, min_pt, *, bold=False):
    """工程表ラベルを折り返さず、幅へ収まるサイズまで縮小する。"""
    size = max_pt
    weight = "bold" if bold else "regular"
    while size >= min_pt - 0.01:
        if text_width_in(text, size, weight) <= width:
            return size
        size -= 0.5
    raise FitError(
        f"{renderer}.{field}: 1行ラベルが最小フォント{min_pt:g}ptでも"
        f"幅{width:.2f}inに収まりません。文言を短くするか期間を広げてください。"
    )


# ---- 番号付きプロセスタイムライン ----
def s_process(slide, spec, page):
    if "flow" in spec:
        return _s_process_flow(slide, spec, page)
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    if spec.get("note") and area.shifted:
        area = ContentArea(area.top, area.bottom - 0.30, area.shifted)
    y = area.map_y
    steps = spec["steps"]
    n = len(steps)
    if not 3 <= n <= 6:
        raise FitError(
            "process: 工程は3〜6件までです。工程を統合または分割してください。")
    left, usable_w = 0.78, 11.72
    w = usable_w / n
    line_y = y(2.42)
    add_rect(slide, left + w / 2, line_y, usable_w - w, 0.035, RULE)
    for i, st in enumerate(steps):
        x = left + i * w
        cx = x + w / 2
        emph = i in spec.get("emph", [])
        color = ACCENT if emph else NAVY
        sp = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - 0.31), Inches(line_y - 0.29),
            Inches(0.62), Inches(0.62))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        add_text(slide, cx - 0.31, line_y - 0.2, 0.62, 0.32, f"{i + 1:02d}",
                 11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        name_size, name_lines = fit_text_or_raise(
            "process", f"steps[{i}].name", st["name"], w - 0.2, 0.42,
            15.5, min_pt=12.5, weight="bold", spacing=1.1)
        add_text(slide, x + 0.1, y(2.94), w - 0.2, 0.42,
                 st["name"], name_size,
                 bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        size, desc_lines = fit_text_or_raise(
            "process", f"steps[{i}].desc", st["desc"],
            w - 0.34, 1.35, 13.5, min_pt=11.5, spacing=1.2)
        add_text(slide, x + 0.17, y(3.56), w - 0.34, 1.35,
                 st["desc"], size,
                 color=TEXT, align=PP_ALIGN.CENTER, spacing=1.2)
        add_text(slide, x + 0.15, y(5.27), w - 0.3, 0.22, "OWNER", 9.5,
                 bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        actor_size, actor_lines = fit_text_or_raise(
            "process", f"steps[{i}].actor", st["actor"],
            w - 0.3, 0.3, 11.5, min_pt=9.5, weight="bold", spacing=1.1)
        add_text(slide, x + 0.15, y(5.56), w - 0.3, 0.3,
                 st["actor"], actor_size,
                 bold=True, color=color, align=PP_ALIGN.CENTER)
    ensure_within(
        "process", y(5.86) - area.top, area.height,
        guidance="工程説明を短くしてください。")
    if spec.get("note"):
        note_line(slide, spec["note"])


def _fit_process_flow(area_h, levels):
    max_rows = max(len(level) for level in levels)

    def candidates():
        values = [
            ("standard", 1.02, 0.34, 12.5),
            ("gap", 1.02, 0.22, 12.5),
            ("element", 0.88, 0.18, 11.0),
        ]
        for stage, node_h, gap_y, font in values:
            used = max_rows * node_h + max(0, max_rows - 1) * gap_y + 0.42
            yield stage, {"node_h": node_h, "gap_y": gap_y, "font": font}, used

    return select_fit(
        "process.flow", area_h, candidates(),
        guidance="工程ノードを減らすか、分岐フローを複数スライドへ分割してください。",
    )


def _s_process_flow(slide, spec, page):
    """levelsとedgesから分岐・合流・戻り線を持つ工程図を描画する。"""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    if spec.get("note") and area.shifted:
        area = ContentArea(area.top, area.bottom - 0.30, area.shifted)
    flow = spec["flow"]
    nodes, levels, edges = flow["nodes"], flow["levels"], flow["edges"]
    fitted = _fit_process_flow(area.height, levels)
    node_h = fitted.values["node_h"]
    gap_y = fitted.values["gap_y"]
    body_size = fitted.values["font"]
    left, usable_w = 0.72, 11.90
    level_count = len(levels)
    # 分岐ラベルを箱に重ねず置けるよう、列数に応じて必要な間隔を確保する。
    gap_x = 0.56 if level_count <= 5 else 0.52
    node_w = (usable_w - gap_x * (level_count - 1)) / level_count
    if node_w < 1.55:
        raise FitError(
            f"process.flow: {level_count}列では箱幅が{node_w:.2f}inとなり"
            "最小1.55inを下回ります。工程階層を減らしてください。")

    rects = {}
    for level_index, level in enumerate(levels):
        top = area.top + 0.34
        x = left + level_index * (node_w + gap_x)
        for row_index, node_id in enumerate(level):
            y = top + row_index * (node_h + gap_y)
            rects[node_id] = (x, y, node_w, node_h)

    feedback_count = 0
    routed = []
    max_node_bottom = max(y + h for _x, y, _w, h in rects.values())
    level_of = {node_id: index for index, level in enumerate(levels)
                for node_id in level}
    for edge in edges:
        source, target = edge["from"], edge["to"]
        sx, sy, sw, sh = rects[source]
        tx, ty, tw, th = rects[target]
        if level_of[target] > level_of[source]:
            start = (sx + sw, sy + sh / 2)
            end = (tx, ty + th / 2)
            trunk = (start[0] + end[0]) / 2
            points = [start, (trunk, start[1]), (trunk, end[1]), end]
            if abs(start[1] - end[1]) < 0.03:
                points = [start, end]
        else:
            lane_y = max_node_bottom + 0.48 + feedback_count * 0.16
            if lane_y > area.bottom - 0.16:
                raise FitError(
                    "process.flow: 差戻し線を本文領域へ配置できません。"
                    "分岐数を減らすか工程を複数スライドへ分割してください。")
            feedback_count += 1
            start = (sx + sw / 2, sy + sh)
            end = (tx + tw / 2, ty + th)
            points = [start, (start[0], lane_y), (end[0], lane_y), end]
        route(slide, points, dash="dash" if edge.get("kind") == "feedback" else None,
              width=1.25)
        routed.append((edge, points))

    for node_id, (x, y, w, h) in rects.items():
        node = nodes[node_id]
        style = node.get("style", "standard")
        fill = LIGHT if style == "decision" else WHITE
        line = ACCENT if style in {"accent", "decision"} else RULE
        add_rect(slide, x, y, w, h, fill, line=line)
        title_size, _ = fit_text_or_raise(
            "process.flow", f"nodes.{node_id}.name", node["name"],
            w - 0.28, 0.32, 14, min_pt=11.5, weight="bold", spacing=1.1)
        add_text(slide, x + 0.14, y + 0.13, w - 0.28, 0.32, node["name"],
                 title_size, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        desc = node.get("desc")
        if desc:
            desc_h = 0.34 if node.get("actor") else h - 0.52
            desc_size, lines = fit_text_or_raise(
                "process.flow", f"nodes.{node_id}.desc", desc,
                w - 0.30, desc_h, body_size, min_pt=9.5, spacing=1.12)
            add_text(slide, x + 0.15, y + 0.50, w - 0.30, desc_h,
                     "\n".join(lines), desc_size, color=TEXT,
                     align=PP_ALIGN.CENTER, spacing=1.12)
        if node.get("actor"):
            actor_size, _ = fit_text_or_raise(
                "process.flow", f"nodes.{node_id}.actor", node["actor"],
                w - 0.30, 0.24, 9.5, min_pt=8.0, weight="bold", spacing=1.05)
            add_text(slide, x + 0.15, y + h - 0.27, w - 0.30, 0.22,
                     node["actor"], actor_size, bold=True, color=ACCENT,
                     align=PP_ALIGN.CENTER)

    for edge, points in routed:
        if not edge.get("label"):
            continue
        if edge.get("kind") == "feedback":
            segments = list(zip(points[:-1], points[1:]))
            start, end = max(
                segments,
                key=lambda segment: abs(segment[1][0] - segment[0][0])
                + abs(segment[1][1] - segment[0][1]),
            )
            label_x, label_y, label_w = (
                (start[0] + end[0]) / 2, (start[1] + end[1]) / 2, 1.35)
        else:
            sx, _sy, sw, _sh = rects[edge["from"]]
            tx, _ty, _tw, _th = rects[edge["to"]]
            end = points[-1]
            label_w = min(1.20, max(0.32, tx - (sx + sw) - 0.06))
            # 分岐元と接続先の間へ置き、どちらの枠も背景マスクで欠かさない。
            label_x = (sx + sw + tx) / 2
            label_y = end[1]
        arrow_label(slide, label_x, label_y, edge["label"],
                    w=label_w, size=8.5)
    if spec.get("note"):
        note_line(slide, spec["note"])


# ---- ロードマップ(フェーズ単位) ----
def s_roadmap(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    months = spec["months"]
    if not 3 <= len(months) <= 12:
        raise FitError(
            "roadmap: 期間は3〜12件までです。期間をまとめるか分割してください。")
    label_x, label_w = 0.72, 2.65
    grid_x, grid_w = label_x + label_w, 9.22
    mw = grid_w / len(months)
    rows = spec["phases"]
    if not 1 <= len(rows) <= 6:
        raise FitError(
            "roadmap: フェーズは1〜6件までです。フェーズを分割してください。")
    fitted = fit_roadmap(area.height, len(rows), has_note=bool(spec.get("note")))
    values = fitted.values
    top = area.top + values["top_gap"]
    hdr_h, row_h = values["header_h"], values["row_h"]
    grid_h = hdr_h + len(rows) * row_h
    add_rect(slide, label_x, top, label_w, hdr_h, NAVY)
    add_text(slide, label_x + 0.18, top + 0.08, label_w - 0.36, 0.28,
             "フェーズ", values["period_pt"], bold=True, color=WHITE)
    add_rect(slide, grid_x, top, grid_w, hdr_h, NAVY)
    for j, m in enumerate(months):
        month_size, _ = fit_text_or_raise(
            "roadmap", f"months[{j}]", m, mw, 0.3, values["period_pt"],
            min_pt=values["period_pt"], weight="bold", spacing=1.1)
        add_text(slide, grid_x + j * mw, top + 0.08, mw, 0.3,
                 m, month_size,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if j:
            add_rect(slide, grid_x + j * mw, top + 0.08, 0.01, hdr_h - 0.16, GRAY)
    for i in range(len(rows)):
        ry = top + hdr_h + i * row_h
        add_rect(slide, label_x, ry, label_w + grid_w, row_h,
                 WHITE if i % 2 == 0 else ZEBRA)
    for j in range(len(months) + 1):
        _grid_line(slide, grid_x + j * mw, top + hdr_h,
                   grid_x + j * mw, top + grid_h)
    for i in range(len(rows) + 1):
        _grid_line(slide, label_x, top + hdr_h + i * row_h,
                   label_x + label_w + grid_w, top + hdr_h + i * row_h)
    for i, ph in enumerate(rows):
        ry = top + hdr_h + i * row_h
        phase_name = re.sub(r"^Phase\s*\d+\s*", "", ph["name"], flags=re.I)
        phase_name = phase_name or ph["name"]
        name_h = min(0.27, row_h * 0.42)
        goal_h = min(0.20, row_h * 0.30)
        add_text(slide, label_x + 0.12, ry + row_h * 0.18, 0.34, name_h,
                 f"{i + 1:02d}", values["name_pt"], bold=True,
                 color=ACCENT)
        phase_size, _ = fit_text_or_raise(
            "roadmap", f"phases[{i}].name", phase_name,
            label_w - 0.62, name_h, values["name_pt"], min_pt=8.5,
            weight="bold", spacing=1.1)
        add_text(slide, label_x + 0.52, ry + row_h * 0.12,
                 label_w - 0.64, name_h,
                 phase_name, phase_size, bold=True, color=NAVY)
        goal_size, _ = fit_text_or_raise(
            "roadmap", f"phases[{i}].goal", ph["goal"],
            label_w - 0.64, goal_h, values["goal_pt"],
            min_pt=7.0, spacing=1.05)
        add_text(slide, label_x + 0.52, ry + row_h * 0.56,
                 label_w - 0.64, goal_h,
                 ph["goal"], goal_size, color=GRAY)
        start, end = resolve_span(ph, months)
        x1 = grid_x + start * mw + 0.05
        x2 = grid_x + end * mw - 0.05
        bar_y = ry + row_h * 0.24
        add_rect(slide, x1, bar_y, x2 - x1, values["bar_h"],
                 ACCENT if i != 1 else NAVY)
        bar_size, _ = fit_text_or_raise(
            "roadmap", f"phases[{i}].bar", ph["bar"],
            x2 - x1 - 0.18, values["bar_h"] - 0.04,
            values["bar_pt"], min_pt=7.5, weight="bold", spacing=1.05)
        add_text(slide, x1 + 0.09, bar_y + 0.035, x2 - x1 - 0.18,
                 values["bar_h"] - 0.04,
                 ph["bar"], bar_size,
                 bold=True, color=WHITE)
    for ms in spec["milestones"]:
        mx = grid_x + resolve_marker(ms["at"], months) * mw
        ry = top + hdr_h + ms["row"] * row_h
        my = ry + min(0.09, row_h * 0.14)
        d = min(0.15, values["bar_h"] * 0.58)
        sp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(mx - d / 2),
                                    Inches(my - d / 2), Inches(d), Inches(d))
        sp.fill.solid()
        sp.fill.fore_color.rgb = CORAL
        sp.line.fill.background()
        sp.shadow.inherit = False
        label_w = min(1.5, max(0.72, mw * 1.6))
        lcx = min(max(mx, grid_x + label_w / 2), grid_x + grid_w - label_w / 2)
        label_size, _ = fit_text_or_raise(
            "roadmap", "milestones.label", ms["label"], label_w,
            min(0.18, row_h * 0.28), values["milestone_pt"],
            min_pt=6.8, spacing=1.0)
        label = add_text(slide, lcx - label_w / 2, ry + row_h * 0.68,
                         label_w, min(0.18, row_h * 0.28), ms["label"],
                         label_size, color=TEXT, align=PP_ALIGN.CENTER,
                         spacing=1.0)
        label.fill.solid()
        label.fill.fore_color.rgb = WHITE if ms["row"] % 2 == 0 else ZEBRA
    if spec.get("note"):
        note_line(slide, spec["note"])


# ---- 複数テーマ・複数作業のプログラム工程表 ----
def s_program_roadmap(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    periods = spec["periods"]
    tracks = spec["tracks"]
    if not 3 <= len(periods) <= 12:
        raise FitError(
            "program_roadmap: 期間は3〜12件までです。期間をまとめてください。")
    if not 1 <= len(tracks) <= 6:
        raise FitError(
            "program_roadmap: テーマは1〜6件までです。複数スライドへ分割してください。")

    packed = [pack_activities(track["activities"], periods) for track in tracks]
    lane_counts = [lane_count for _placements, lane_count in packed]
    fitted = fit_program_roadmap(
        area.height, lane_counts, has_note=bool(spec.get("note")))
    values = fitted.values

    label_x, track_col_w = 0.72, 2.78
    grid_x, grid_w = label_x + track_col_w, 9.09
    period_w = grid_w / len(periods)
    top = area.top + values["top_gap"]
    header_h = values["header_h"]
    row_heights = [
        2 * values["track_pad"] + count * values["lane_pitch"]
        for count in lane_counts
    ]
    grid_h = header_h + sum(row_heights) + values["track_gap"] * (len(tracks) - 1)

    add_rect(slide, label_x, top, track_col_w, header_h, NAVY)
    add_text(slide, label_x + 0.18, top + 0.07, track_col_w - 0.36, 0.28,
             "テーマ", values["period_pt"], bold=True, color=WHITE)
    add_rect(slide, grid_x, top, grid_w, header_h, NAVY)
    for index, period in enumerate(periods):
        size, _ = fit_text_or_raise(
            "program_roadmap", f"periods[{index}]", period,
            period_w - 0.04, 0.28, values["period_pt"],
            min_pt=7.5, weight="bold", spacing=1.0)
        add_text(slide, grid_x + index * period_w, top + 0.07,
                 period_w, 0.28, period, size, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, spacing=1.0)

    row_tops = []
    cursor = top + header_h
    for index, row_h in enumerate(row_heights):
        row_tops.append(cursor)
        add_rect(slide, label_x, cursor, track_col_w + grid_w, row_h,
                 WHITE if index % 2 == 0 else ZEBRA)
        cursor += row_h + values["track_gap"]

    body_top = top + header_h
    for index in range(len(periods) + 1):
        _grid_line(slide, grid_x + index * period_w, body_top,
                   grid_x + index * period_w, top + grid_h)
    for index, row_top in enumerate(row_tops):
        _grid_line(slide, label_x, row_top,
                   label_x + track_col_w + grid_w, row_top)
        track = tracks[index]
        row_h = row_heights[index]
        add_text(slide, label_x + 0.12, row_top + 0.11, 0.38, 0.25,
                 f"{index + 1:02d}", values["track_pt"], bold=True,
                 color=ACCENT)
        track_size = _fit_single_line(
            "program_roadmap", f"tracks[{index}].name", track["name"],
            track_col_w - 0.66, values["track_pt"], 8.5, bold=True)
        add_text(slide, label_x + 0.52, row_top + 0.09,
                 track_col_w - 0.66, min(0.28, row_h - 0.14),
                 track["name"], track_size, bold=True, color=NAVY,
                 spacing=1.0, wrap=False)

        placements, _lane_count = packed[index]
        lane_neighbors = {}
        for lane in range(_lane_count):
            lane_items = sorted(
                (placement for placement in placements if placement.lane == lane),
                key=lambda placement: (placement.start, placement.end),
            )
            for position, placement in enumerate(lane_items):
                previous = lane_items[position - 1] if position else None
                following = lane_items[position + 1] if position + 1 < len(lane_items) else None
                lane_neighbors[placement.index] = (previous, following)
        for placement in placements:
            activity = placement.activity
            lane_top = (row_top + values["track_pad"]
                        + placement.lane * values["lane_pitch"])
            x1 = grid_x + placement.start * period_w + 0.035
            x2 = grid_x + placement.end * period_w - 0.035
            color = CORAL if activity.get("emph") else ACCENT
            _activity_line(slide, x1, lane_top + 0.035, x2, color)
            text_h = values["lane_pitch"] - 0.07
            previous, following = lane_neighbors[placement.index]
            label_center = (x1 + x2) / 2
            previous_center = None if previous is None else (
                grid_x + (previous.start + previous.end) * period_w / 2)
            following_center = None if following is None else (
                grid_x + (following.start + following.end) * period_w / 2)
            activity_label_x, activity_label_w = centered_label_box(
                label_center, previous_center, following_center,
                grid_x, grid_x + grid_w, 1.45,
            )
            activity_size = _fit_single_line(
                "program_roadmap",
                f"tracks[{index}].activities[{placement.index}].label",
                activity["label"], activity_label_w - 0.04,
                values["activity_pt"], 7.5,
                bold=bool(activity.get("emph")))
            label = add_text(
                slide, activity_label_x + 0.02, lane_top + 0.07,
                activity_label_w - 0.04, text_h, activity["label"], activity_size,
                bold=bool(activity.get("emph")), color=TEXT,
                align=PP_ALIGN.CENTER, spacing=1.0, wrap=False)
            label.fill.solid()
            label.fill.fore_color.rgb = WHITE if index % 2 == 0 else ZEBRA
    _grid_line(slide, label_x, cursor - values["track_gap"],
               label_x + track_col_w + grid_w, cursor - values["track_gap"])
    if spec.get("note"):
        note_line(slide, spec["note"])


# ---- 2軸ポジショニングマップ ----
def s_matrix(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    if spec.get("note") and area.shifted:
        area = ContentArea(area.top, area.bottom - 0.30, area.shifted)
    y = area.map_y
    if not 1 <= len(spec["points"]) <= 8:
        raise FitError(
            "matrix: 点は1〜8件までです。点をまとめるかスライドを分割してください。")
    ox, oy = 1.68, y(6.24)
    aw, ah = 10.5, oy - y(2.22)
    mid_x, mid_y = ox + aw / 2, oy - ah / 2

    quadrants = spec.get("quadrants")
    if quadrants:
        # Explicit region names make the highlighted quadrant semantic. Existing
        # matrix specs without them retain the neutral scatter-plot treatment.
        add_rect(slide, ox, oy - ah, aw, ah, WHITE, line=RULE)
        add_rect(slide, mid_x, oy - ah, aw / 2, ah / 2, LIGHT)
        qlabels = [
            (ox + 0.18, oy - 0.42, 2.2, quadrants[0], PP_ALIGN.LEFT, GRAY),
            (ox + aw - 2.38, oy - 0.42, 2.2, quadrants[1], PP_ALIGN.RIGHT, GRAY),
            (ox + 0.18, oy - ah + 0.18, 2.2, quadrants[2], PP_ALIGN.LEFT, GRAY),
            (ox + aw - 2.38, oy - ah + 0.18, 2.2, quadrants[3], PP_ALIGN.RIGHT, ACCENT),
        ]
        for idx, (x, y, w, label, align, color) in enumerate(qlabels):
            label_size, _ = fit_text_or_raise(
                "matrix", f"quadrants[{idx}]", label, w, 0.3, 10.5,
                min_pt=8.5, weight="bold", spacing=1.1)
            add_text(slide, x, y, w, 0.3, label, label_size, bold=True,
                     color=color, align=align)
    else:
        target_size, _ = fit_text_or_raise(
            "matrix", "target_label", spec["target_label"], 2.0, 0.3, 10.5,
            min_pt=8.5, weight="bold", spacing=1.1)
        add_text(slide, mid_x + 0.18, oy - ah + 0.16, 2.0, 0.3,
                 spec["target_label"], target_size, bold=True, color=ACCENT)

    add_arrow(slide, ox, oy, ox + aw, oy, width=1.75)
    add_arrow(slide, ox, oy, ox, oy - ah, width=1.75)
    x_axis_size, _ = fit_text_or_raise(
        "matrix", "x_axis", spec["x_axis"], 2.4, 0.3, 11.5,
        min_pt=9.5, weight="bold", spacing=1.1)
    y_axis_size, _ = fit_text_or_raise(
        "matrix", "y_axis", spec["y_axis"], 3.0, 0.3, 11.5,
        min_pt=9.5, weight="bold", spacing=1.1)
    add_text(slide, ox + aw - 2.5, oy + 0.16, 2.4, 0.3, spec["x_axis"], x_axis_size,
             bold=True, color=TEXT, align=PP_ALIGN.RIGHT)
    add_text(slide, ox - 0.02, oy - ah - 0.38, 3.0, 0.3, spec["y_axis"], y_axis_size,
             bold=True, color=TEXT)
    _stage, label_w, point_xy, labels = fit_matrix_labels(
        spec["points"], ox, oy - ah, aw, ah,
        barriers=(mid_x, mid_y) if quadrants else ())
    for p, (px, py) in zip(spec["points"], point_xy):
        r = 0.15
        sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - r), Inches(py - r),
                                    Inches(2 * r), Inches(2 * r))
        sp.fill.solid()
        sp.fill.fore_color.rgb = CORAL if p.get("emph") else ACCENT
        sp.line.color.rgb = WHITE
        sp.line.width = Pt(1.0)
        sp.shadow.inherit = False
    for index, p in enumerate(spec["points"]):
        left, top, right, bottom = labels[index]
        point_size, point_lines = fit_text_or_raise(
            "matrix", "points.name", p["name"], label_w, 0.3, 11,
            min_pt=9, weight="bold" if p.get("emph") else "regular",
            spacing=1.1)
        add_text(slide, left, top, right - left, bottom - top,
                 p["name"], point_size,
                 bold=bool(p.get("emph")), color=CORAL if p.get("emph") else TEXT,
                 align=PP_ALIGN.CENTER)
    ensure_within(
        "matrix", oy - area.top, area.height,
        guidance="軸ラベルまたは点の名称を短くしてください。")
    if spec.get("note"):
        note_line(slide, spec["note"])
