"""大判画像rendererのパス・比率・収容段階を検証する。"""
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

import generate
from asset_paths import ASSET_DIR, resolve_image_path
from image_slide import PICTURE_NAME, fit_image_layout, s_image
from layout_fit import FitError
from validate_content import validate


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _must_fail(fn, expected):
    try:
        fn()
    except FitError as exc:
        assert expected in str(exc), str(exc)
    else:
        raise AssertionError("過密入力を拒否しませんでした")


def main():
    assert resolve_image_path("cover/cover-background.png") == (
        ASSET_DIR / "cover" / "cover-background.png").resolve()

    spec = {
        "type": "image", "kicker": "画像", "title": "大判画像",
        "image": "cover/cover-background.png", "fit": "contain",
        "alt": "濃紺を基調とした抽象的な背景画像",
    }
    deck = {
        "meta": {"title": "検証", "footer": "検証", "date": "2026年7月",
                 "author": "検証担当"},
        "slides": [spec],
    }
    assert not validate(deck)

    contain_slide = _slide()
    s_image(contain_slide, spec, 1)
    contain = next(shape for shape in contain_slide.shapes
                   if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert contain.name == PICTURE_NAME
    assert contain.crop_left == contain.crop_right == 0
    assert contain.crop_top == contain.crop_bottom == 0
    assert contain._element.nvPicPr.cNvPr.get("descr") == spec["alt"]
    assert contain._element.spPr.find(qn("a:effectLst")) is None

    cover_spec = dict(spec, fit="cover")
    cover_slide = _slide()
    s_image(cover_slide, cover_spec, 1)
    cover = next(shape for shape in cover_slide.shapes
                 if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert cover.crop_top > 0 or cover.crop_left > 0

    shadow_slide = _slide()
    s_image(shadow_slide, dict(spec, shadow=True), 1)
    shadow_picture = next(
        shape for shape in shadow_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    effect_list = shadow_picture._element.spPr.find(qn("a:effectLst"))
    shadow = effect_list.find(qn("a:outerShdw"))
    assert shadow.get("dir") == "2700000"
    assert shadow.get("dist") == "50800"

    assert fit_image_layout(5.27).stage == "standard"
    assert fit_image_layout(3.50).stage == "gap"
    assert fit_image_layout(3.15).stage == "element"
    _must_fail(lambda: fit_image_layout(2.40), "最小設定")

    invalid = deepcopy(deck)
    invalid["slides"][0]["image"] = "../../outside.png"
    assert any("assets内" in error for error in validate(invalid))
    invalid = deepcopy(deck)
    invalid["slides"][0]["image"] = "images/missing.png"
    assert any("assets/にありません" in error for error in validate(invalid))
    invalid = deepcopy(deck)
    invalid["slides"][0]["fit"] = "stretch"
    assert any("contain" in error for error in validate(invalid))
    invalid = deepcopy(deck)
    invalid["slides"][0]["shadow"] = "yes"
    assert any("true または false" in error for error in validate(invalid))
    invalid = deepcopy(deck)
    invalid["slides"][0]["caption"] = "説明"
    invalid["slides"][0]["source"] = "出典"
    errors = validate(invalid)
    assert any("caption" in error and "lead" in error for error in errors)
    assert any("source" in error and "削除" in error for error in errors)

    print("image slide tests passed")


if __name__ == "__main__":
    main()
