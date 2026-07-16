"""公開diagramスキーマの構成図サンプルを回帰検証する。"""
from copy import deepcopy
import json

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from diagram_layout import ICON_SIZE, Layout, render_diagram
from diagram_specs import AWS_MULTIAZ_EXAMPLE, AWS_SIMPLE_EXAMPLE
from diagrams import EDGE_GAP, NODE_LABEL_PAD_X, NODE_LABEL_PAD_Y
from layout_fit import FitError
from textfit import line_height_in, text_width_in
from validate_content import validate


def _assert_no_layout_fields(diagram):
    assert "area" not in diagram
    for container in diagram.get("containers", []):
        assert "pad" not in container
        assert "pad_x" not in container


def main():
    examples = [AWS_SIMPLE_EXAMPLE, AWS_MULTIAZ_EXAMPLE]
    for diagram in examples:
        _assert_no_layout_fields(diagram)
        layout = Layout(diagram)
        edges = diagram["edges"]
        routed = layout.route_edges(edges)
        layout.validate_edges(edges, routed)

    # 上下ポートはラベル下端ではなくアイコン外周に置き、線をラベル背面へ通す。
    simple_layout = Layout(AWS_SIMPLE_EXAMPLE)
    node_name = next(name for name, node in AWS_SIMPLE_EXAMPLE["nodes"].items()
                     if node.get("sub"))
    cx, cy = simple_layout.node_center(node_name)
    bottom = simple_layout.port(node_name, "bottom")
    assert abs(bottom[0] - cx) <= 0.001
    assert abs(bottom[1] - (cy + simple_layout.icon_r + EDGE_GAP)) <= 0.001
    assert bottom[1] < simple_layout._label_bottom(node_name)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_diagram(slide, AWS_SIMPLE_EXAMPLE)
    shapes = list(slide.shapes)
    line_z = [i for i, shape in enumerate(shapes)
              if shape.shape_type == MSO_SHAPE_TYPE.LINE]
    picture_z = [i for i, shape in enumerate(shapes)
                 if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert line_z and picture_z
    assert max(line_z) < min(picture_z), "コネクタはノード画像より先に描画する"

    expected = {
        value
        for node in AWS_SIMPLE_EXAMPLE["nodes"].values()
        for value in (node["title"], node.get("sub"))
        if value
    }
    masked = [shape for shape in shapes
              if shape.has_text_frame and shape.text_frame.text in expected]
    assert len(masked) == len(expected)
    for shape in masked:
        text = shape.text_frame.text
        run = shape.text_frame.paragraphs[0].runs[0]
        size = run.font.size.pt
        bold = bool(run.font.bold)
        weight = "bold" if bold else "regular"
        assert str(shape.fill.type) == "SOLID (1)"
        assert shape.width.inches <= text_width_in(text, size, weight) \
            + NODE_LABEL_PAD_X + 0.01
        assert shape.height.inches <= line_height_in(size, 1.1) \
            + NODE_LABEL_PAD_Y + 0.01

    dense_slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_diagram(dense_slide, AWS_MULTIAZ_EXAMPLE)
    name_label = next(shape for shape in dense_slide.shapes
                      if shape.has_text_frame
                      and shape.text_frame.text == "名前解決")
    dense_layout = Layout(AWS_MULTIAZ_EXAMPLE)
    r53_x, _ = dense_layout.node_center("r53")
    _, cf_y = dense_layout.node_center("cf")
    assert name_label.left.inches + name_label.width.inches < r53_x
    assert name_label.top.inches + name_label.height.inches < cf_y - 0.10

    deck = {
        "meta": {
            "title": "diagram examples",
            "footer": "diagram examples",
            "date": "2026-07",
            "author": "test",
        },
        "slides": [
            {
                "type": "diagram",
                "kicker": "simple",
                "title": "simple architecture",
                "diagram": AWS_SIMPLE_EXAMPLE,
            },
            {
                "type": "diagram",
                "kicker": "dense",
                "title": "dense architecture",
                "diagram": AWS_MULTIAZ_EXAMPLE,
            },
        ],
    }
    json_roundtrip = json.loads(json.dumps(deck, ensure_ascii=False))
    errors = validate(json_roundtrip)
    assert not errors, "\n".join(errors)

    invalid = deepcopy(json_roundtrip)
    invalid["slides"][0]["diagram"]["area"] = [0, 0, 1, 1]
    invalid["slides"][0]["diagram"]["containers"][0]["pad"] = 0.1
    invalid["slides"][0]["diagram"]["nodes"]["user"]["icon"] = (
        "fluent/not_defined.png")
    errors = validate(invalid)
    assert any("diagram.area" in error for error in errors)
    assert any(".pad" in error for error in errors)
    assert any("assets/ にありません" in error for error in errors)

    rows = [f"r{i}" for i in range(5)]
    nodes = {
        f"n{i}": {
            "col": "main", "row": row, "icon": "fluent/app.png",
            "title": f"Node {i}", **({"sub": "sub"} if i < 2 else {}),
        }
        for i, row in enumerate(rows)
    }
    dense = {
        "cols": ["main"], "rows": rows, "nodes": nodes,
        "containers": [], "channels": {}, "edges": [],
    }
    gap_only = deepcopy(dense)
    gap_only["nodes"]["n0"].pop("sub")
    gap_only["nodes"]["n1"].pop("sub")
    compact_gap = Layout(gap_only)
    assert compact_gap.fit_stage == "gap"
    assert compact_gap.icon_size == ICON_SIZE

    compact = Layout(dense)
    assert compact.fit_stage == "icon"
    assert compact.icon_size < ICON_SIZE

    for node in dense["nodes"].values():
        node["sub"] = "sub"
    try:
        Layout(dense)
    except FitError as e:
        assert "アイコン縮小" in str(e), str(e)
        assert "行数を減らす" in str(e), str(e)
    else:
        raise AssertionError("最小アイコンでも収まらないdiagramを拒否しませんでした")
    print("OK: diagram examples passed schema, layout, routing, and label-mask checks")


if __name__ == "__main__":
    main()
