"""python-pptxとテキスト実測による基本サンプルデッキ生成。"""
import argparse
import math
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from content import DECK
from cover_footer import load_cover_footer_config, render_cover, render_footer
from layout_fit import FitError, fit_text_or_raise, select_fit, stepped
from textfit import line_height_in, text_width_in, wrap_text

# ---- テーマ ----
NAVY = RGBColor(0x18, 0x2C, 0x43)
ACCENT = RGBColor(0x0D, 0x78, 0x70)
CORAL = RGBColor(0xC7, 0x58, 0x3E)
LIGHT = RGBColor(0xDF, 0xEB, 0xE8)
TEXT = RGBColor(0x20, 0x27, 0x29)
GRAY = RGBColor(0x66, 0x6E, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFC)
ZEBRA = RGBColor(0xEC, 0xEA, 0xE4)
CANVAS = RGBColor(0xF7, 0xF5, 0xEF)
RULE = RGBColor(0xD1, 0xCF, 0xC8)
FONT = "Yu Gothic"
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.55
BODY_W = SLIDE_W - MARGIN * 2
BODY_TOP, BODY_BOTTOM = 1.58, 6.85
LEAD_Y, LEAD_MAX_H = 1.58, 0.56
TABLE_HEADER_H = 0.72
TABLE_TOP_GAP = 0.38
TABLE_BOTTOM_GAP = 0.15
TABLE_NOTE_H = 0.30
COVER_FOOTER = load_cover_footer_config()


@dataclass(frozen=True)
class ContentArea:
    """ヘッダー下のrenderer描画領域。leadなしでは従来値を保持する。"""

    top: float = BODY_TOP
    bottom: float = BODY_BOTTOM
    shifted: bool = False

    @property
    def height(self):
        return self.bottom - self.top

    def map_y(self, y):
        """固定構図の従来Y座標を、lead指定時の本文領域へ写像する。"""
        if not self.shifted:
            return y
        scale = self.height / (BODY_BOTTOM - BODY_TOP)
        return self.top + (y - BODY_TOP) * scale


def configure_cover_footer(path=None):
    """表紙・フッター設定を切り替える。未指定なら標準設定へ戻す。"""
    global COVER_FOOTER
    COVER_FOOTER = load_cover_footer_config(path)


def set_run(run, size, *, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "ja-JP")
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def add_text(slide, x, y, w, h, text, size, *, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.3,
             wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        set_run(p.add_run(), size, bold=bold, color=color)
        p.runs[0].text = line
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None, round_=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def header(slide, kicker, title, lead=None):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CANVAS)
    kicker_size, _ = fit_text_or_raise(
        "header", "kicker", kicker, 4.8, 0.32, 11.5,
        min_pt=9, weight="bold", spacing=1.1)
    add_text(slide, 0.72, 0.27, 4.8, 0.32, kicker, kicker_size,
             bold=True, color=ACCENT)
    size = 27
    lines = wrap_text(title, 11.9, size, "bold")
    while len(lines) > 1 and size > 18:
        size -= 0.5
        lines = wrap_text(title, 11.9, size, "bold")
    if len(lines) > 1:
        size, lines = fit_text_or_raise(
            "header", "title", title, 11.9, 0.86, 18,
            min_pt=16, weight="bold", spacing=1.12)
    add_text(slide, 0.72, 0.67, 11.9, 0.86, "\n".join(lines), size,
             bold=True, color=NAVY, spacing=1.12)
    if not lead:
        return ContentArea()

    lead_size, lead_lines = fit_text_or_raise(
        "header", "lead", lead, 11.9, LEAD_MAX_H, 14,
        min_pt=11.5, spacing=1.18)
    lead_h = len(lead_lines) * line_height_in(lead_size, 1.18) + 0.03
    add_text(slide, 0.72, LEAD_Y, 11.9, lead_h, "\n".join(lead_lines),
             lead_size, color=GRAY, spacing=1.18)
    return ContentArea(top=LEAD_Y + lead_h + 0.17, shifted=True)


def page_label(page):
    """Return a stable page marker shared by every generator entry point."""
    total = len(DECK["slides"])
    digits = max(2, len(str(total)))
    return f"{page:0{digits}d} / {total:0{digits}d}"


def footer(slide, page):
    total = len(DECK["slides"])
    render_footer(slide, page, DECK["meta"], total, COVER_FOOTER,
                  add_text=add_text, add_rect=add_rect)


def note_line(slide, note):
    size, _ = fit_text_or_raise(
        "note", "text", note, BODY_W, 0.25, 8.5,
        min_pt=7, spacing=1.1)
    add_text(slide, MARGIN, 6.62, BODY_W, 0.25, note, size,
             color=GRAY, align=PP_ALIGN.RIGHT)


# ---- スライド種別 ----
def s_title(slide, spec, page):
    meta = DECK["meta"]
    total = len(DECK["slides"])
    render_cover(slide, spec, meta, total, COVER_FOOTER,
                 add_text=add_text, add_rect=add_rect)


def s_bullets(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    bullets = spec["bullets"]
    area_h = area.height - 0.22
    tx = 1.68
    tw = 10.15
    def measure(size, gap):
        heights = [len(wrap_text(t, tw, size)) * line_height_in(size, 1.22)
                   for t, _ in bullets]
        total = sum(heights) + gap * (len(bullets) - 1)
        return heights, total

    def candidates():
        for gap in stepped(0.48, 0.30, 0.03):
            _heights, total = measure(18, gap)
            yield ("standard" if gap == 0.48 else "gap",
                   {"size": 18, "gap": gap}, total)
        for size in stepped(17.5, 12, 0.5):
            _heights, total = measure(size, 0.30)
            yield "font", {"size": size, "gap": 0.30}, total

    fitted = select_fit(
        "bullets", area_h, candidates(),
        guidance="箇条書きを減らすか各項目を短くしてください。",
    )
    size, gap = fitted.values["size"], fitted.values["gap"]
    heights, total = measure(size, gap)
    y = area.top + 0.38 + max(0.0, (area_h - total) * 0.22)
    for i, ((text, _), bh) in enumerate(zip(bullets, heights), 1):
        add_text(slide, 0.78, y - 0.07, 0.62, 0.45, f"{i:02d}", 15,
                 bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
        add_text(slide, tx, y, tw, bh + 0.08, text, size, spacing=1.22)
        if i < len(bullets):
            add_rect(slide, tx, y + bh + 0.15, tw, 0.012, RULE)
        y += bh + gap


def s_cards(slide, spec, page):
    """Render purpose-specific cards instead of one universal panel pattern."""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    cards = [_normalize_card(card) for card in spec["cards"]]
    n = len(cards)
    style = spec.get("style", "editorial")
    left, usable_w = 0.78, 11.78

    if style == "metrics":
        cols = n if n <= 4 else 3
        rows = math.ceil(n / cols)
        gap_x, gap_y = 0.46, 0.42
        cw = (usable_w - gap_x * (cols - 1)) / cols
        ch = min(2.72, (area.height - 0.58 - gap_y * (rows - 1)) / rows)
        if ch < 1.62:
            raise FitError(
                "cards.metrics: KPIカードの高さが不足しています。カード数または本文を"
                "減らしてください。")
        top = area.top + 0.42
        body_size = min(
            fit_text_or_raise(
                "cards.metrics", f"cards[{i}].body", card["body"],
                cw - 0.16, ch - 1.26, 13, min_pt=10.5, spacing=1.16,
            )[0]
            for i, card in enumerate(cards))
        for i, card in enumerate(cards):
            row, col = divmod(i, cols)
            x = left + col * (cw + gap_x)
            y = top + row * (ch + gap_y)
            label, value = card["heading"], card.get("value", "")
            if not value:
                label, value = _split_metric_head(label)
            label_size, _ = fit_text_or_raise(
                "cards.metrics", f"cards[{i}].label", label,
                cw, 0.34, 12.5, min_pt=10.5, weight="bold", spacing=1.1)
            add_text(slide, x, y, cw, 0.34, label, label_size,
                     bold=True, color=NAVY)
            value_size = fit_text_or_raise(
                "cards.metrics", f"cards[{i}].value", value,
                cw, 0.66, 32, min_pt=22, weight="bold")[0]
            color = CORAL if card["emphasis"] else ACCENT
            add_text(slide, x, y + 0.42, cw, 0.66, value, value_size,
                     bold=True, color=color)
            add_text(slide, x, y + 1.22, cw, ch - 1.22, card["body"],
                     body_size, color=TEXT, spacing=1.2)
            if col < cols - 1 and i + 1 < n:
                add_rect(slide, x + cw + gap_x / 2, y + 0.04,
                         0.012, ch - 0.08, RULE)
            if row < rows - 1:
                add_rect(slide, x, y + ch + gap_y / 2, cw, 0.012, RULE)
        return

    emphasized = [i for i, card in enumerate(cards) if card["emphasis"]]
    if style == "editorial" and n == 4 and len(emphasized) == 1:
        lead_index = emphasized[0]
        lead_card = cards[lead_index]
        supporting = [card for i, card in enumerate(cards) if i != lead_index]
        top = area.top + 0.46
        add_text(slide, 0.8, top, 0.58, 0.36, "01", 15,
                 bold=True, color=GRAY)
        lead_head_size, _ = fit_text_or_raise(
            "cards.editorial", f"cards[{lead_index}].heading", lead_card["heading"],
            3.95, 0.56, 23, min_pt=18, weight="bold", spacing=1.1)
        add_text(slide, 1.52, top - 0.03, 3.95, 0.56, lead_card["heading"],
                 lead_head_size, bold=True, color=NAVY)
        lead_size, lead_lines = fit_text_or_raise(
            "cards.editorial", f"cards[{lead_index}].body", lead_card["body"],
            4.1, 2.05, 16, min_pt=13, spacing=1.28)
        add_text(slide, 1.52, top + 0.76, 4.1, 2.05,
                 "\n".join(lead_lines), lead_size, color=TEXT, spacing=1.28)

        right_x, right_w = 6.42, 5.92
        row_h = 1.31
        for i, card in enumerate(supporting, 2):
            y = top + (i - 2) * row_h
            add_text(slide, right_x, y + 0.02, 0.5, 0.32, f"{i:02d}", 12.5,
                     bold=True, color=GRAY)
            head_size, _ = fit_text_or_raise(
                "cards.editorial", f"supporting[{i - 2}].heading", card["heading"],
                right_w - 0.68, 0.38, 16, min_pt=13,
                weight="bold", spacing=1.1)
            add_text(slide, right_x + 0.68, y, right_w - 0.68, 0.38,
                     card["heading"],
                     head_size, bold=True, color=NAVY)
            body_size, body_lines = fit_text_or_raise(
                "cards.editorial", f"supporting[{i - 2}].body", card["body"],
                right_w - 0.72, 0.56, 12.5, min_pt=11, spacing=1.15)
            add_text(slide, right_x + 0.68, y + 0.48, right_w - 0.72, 0.56,
                     "\n".join(body_lines), body_size,
                     color=TEXT, spacing=1.15)
            if i < n:
                add_rect(slide, right_x + 0.68, y + 1.14, right_w - 0.68, 0.012, RULE)
        return

    cols = n if n <= 3 else (2 if n == 4 else 3)
    rows = math.ceil(n / cols)
    gap_x, gap_y = 0.72, 0.44
    cw = (usable_w - gap_x * (cols - 1)) / cols
    area_h = area.height - 0.62
    fit_available = area_h if rows == 2 else min(3.32, area_h)

    def measure(size, candidate_gap):
        row_needs = []
        for row in range(rows):
            row_cards = cards[row * cols:(row + 1) * cols]
            body_need = max(
                len(wrap_text(card["body"], cw - 0.92, size))
                * line_height_in(size, 1.2)
                for card in row_cards)
            row_needs.append(0.72 + body_need)
        return sum(row_needs) + candidate_gap * (rows - 1)

    def candidates():
        gaps = stepped(gap_y, 0.28, 0.04) if rows == 2 else [gap_y]
        for candidate_gap in gaps:
            used = measure(14, candidate_gap)
            yield ("standard" if candidate_gap == gap_y else "gap",
                   {"size": 14, "gap": candidate_gap}, used)
        for size in stepped(13.5, 11, 0.5):
            used = measure(size, 0.28 if rows == 2 else gap_y)
            yield "font", {"size": size,
                            "gap": 0.28 if rows == 2 else gap_y}, used

    fitted = select_fit(
        "cards", fit_available, candidates(),
        guidance="カード本文を短くするかカード数を減らしてください。",
    )
    body_size = fitted.values["size"]
    gap_y = fitted.values["gap"]
    ch = ((area_h - gap_y * (rows - 1)) / rows
          if rows > 1 else min(3.32, area_h))
    top = area.top + (0.42 if rows == 2 else 0.72)
    for i, card in enumerate(cards):
        row, col = divmod(i, cols)
        x = left + col * (cw + gap_x)
        y = top + row * (ch + gap_y)
        add_text(slide, x, y + 0.02, 0.58, 0.36, f"{i + 1:02d}", 14,
                 bold=True, color=GRAY)
        head_size, _ = fit_text_or_raise(
            "cards", f"cards[{i}].heading", card["heading"],
            cw - 0.72, 0.42, 16.5, min_pt=13,
            weight="bold", spacing=1.1)
        head_color = CORAL if card["emphasis"] else NAVY
        add_text(slide, x + 0.72, y, cw - 0.72, 0.42,
                 card["heading"], head_size, bold=True, color=head_color)
        add_text(slide, x + 0.72, y + 0.58, cw - 0.82, ch - 0.64, card["body"],
                 body_size, color=TEXT, spacing=1.2)
        if row < rows - 1:
            add_rect(slide, x + 0.72, y + ch + gap_y / 2, cw - 0.72, 0.012, RULE)


def _split_metric_head(head):
    match = re.match(r"^(.*?)[\s　]+([+\-−]?[0-9][0-9,.]*\s*(?:%|分|件|倍|pt)?)$", head)
    if match:
        return match.group(1).strip(), match.group(2).replace("−", "-").strip()
    return head, ""


def _normalize_card(card):
    """旧2要素配列と新しい意味付きobjectを共通形式へ変換する。"""
    if isinstance(card, dict):
        return {
            "heading": card["heading"],
            "body": card["body"],
            "value": card.get("value", ""),
            "emphasis": bool(card.get("emphasis")),
        }
    return {
        "heading": card[0], "body": card[1], "value": "", "emphasis": False,
    }


def _auto_table_widths(columns, rows):
    """見出しとセル内容を実測し、本文幅を列へ決定論的に配分する。"""
    count = len(columns)
    min_width = max(0.92, min(1.38, BODY_W / count * 0.58))
    natural = []
    for index, heading in enumerate(columns):
        samples = [text_width_in(heading, 12, "bold")]
        samples.extend(text_width_in(row[index], 11.5) for row in rows)
        natural.append(min(4.65, max(min_width, max(samples) + 0.34)))

    widths = [min_width] * count
    remaining = BODY_W - sum(widths)
    unmet = {index for index in range(count) if natural[index] > widths[index]}
    while remaining > 0.001 and unmet:
        share = remaining / len(unmet)
        used = 0.0
        completed = set()
        for index in unmet:
            addition = min(share, natural[index] - widths[index])
            widths[index] += addition
            used += addition
            if natural[index] - widths[index] <= 0.001:
                completed.add(index)
        remaining -= used
        unmet -= completed
        if used <= 0.001:
            break
    if remaining > 0:
        weights = [max(1.0, value) for value in natural]
        total_weight = sum(weights)
        widths = [width + remaining * weight / total_weight
                  for width, weight in zip(widths, weights)]
    widths[-1] += BODY_W - sum(widths)
    return widths


def _fit_table(rows, widths, avail):
    """表の余白圧縮・フォント縮小を選び、行高と判定結果を返す。"""
    min_row_h = min(1.04, max(0.64, avail / max(1, len(rows))))

    def measure(size, pad):
        row_hs = []
        for row in rows:
            need = max(
                len(wrap_text(c, widths[j] - pad * 2, size))
                * line_height_in(size, 1.15) for j, c in enumerate(row))
            row_hs.append(max(need + pad * 2, min_row_h))
        return row_hs, sum(row_hs)

    def candidates():
        for pad in stepped(0.14, 0.10, 0.02):
            _row_hs, used = measure(13.5, pad)
            yield ("standard" if pad == 0.14 else "padding",
                   {"size": 13.5, "pad": pad}, used)
        for size in stepped(13.0, 10.5, 0.5):
            _row_hs, used = measure(size, 0.10)
            yield "font", {"size": size, "pad": 0.10}, used

    fitted = select_fit(
        "table", avail, candidates(),
        guidance="表の行を減らすかセル内の文言を短くしてください。",
    )
    row_hs, _used = measure(fitted.values["size"], fitted.values["pad"])
    return fitted, row_hs


def s_table(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    cols, rows = spec["columns"], spec["rows"]
    widths = spec.get("col_widths") or _auto_table_widths(cols, rows)
    assert abs(sum(widths) - BODY_W) < 0.6, f"列幅合計={sum(widths)}"
    hdr_h = TABLE_HEADER_H
    table_available = (
        area.height - TABLE_TOP_GAP - TABLE_BOTTOM_GAP
        - (TABLE_NOTE_H if spec.get("note") else 0)
    )
    avail = table_available - hdr_h
    fitted, row_hs = _fit_table(rows, widths, avail)
    size, pad = fitted.values["size"], fitted.values["pad"]
    table_h = hdr_h + sum(row_hs)
    top = area.top + TABLE_TOP_GAP + max(0.0, (table_available - table_h) * 0.12)
    gt = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(MARGIN),
                                Inches(top), Inches(BODY_W), Inches(table_h))
    table = gt.table
    table.first_row = False
    table.horz_banding = False
    for j, wdt in enumerate(widths):
        table.columns[j].width = Emu(int(Inches(wdt)))
    table.rows[0].height = Emu(int(Inches(hdr_h)))
    for i, rh in enumerate(row_hs):
        table.rows[i + 1].height = Emu(int(Inches(rh)))
    for j, name in enumerate(cols):
        header_size, _ = fit_text_or_raise(
            "table", f"columns[{j}]", name, widths[j] - 0.18, hdr_h - 0.08,
            size, min_pt=10.5, weight="bold", spacing=1.15)
        _cell(table.cell(0, j), name, header_size, bold=True, color=WHITE, fill=NAVY,
              center=j != len(cols) - 1 and j != 0)
    for i, row in enumerate(rows):
        fill = WHITE if i % 2 else ZEBRA
        for j, val in enumerate(row):
            _cell(table.cell(i + 1, j), val, size,
                  bold=(j == 0), color=NAVY if j == 0 else TEXT, fill=fill,
                  center=0 < j < len(cols) - 1 and len(val) <= 6)
    if spec.get("note"):
        note_line(slide, spec["note"])


def _cell(cell, text, size, *, bold=False, color=TEXT, fill=WHITE, center=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.09)
    cell.margin_top = cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        p.line_spacing = 1.15
        set_run(p.add_run(), size, bold=bold, color=color)
        p.runs[0].text = line


def s_twocol(slide, spec, page):
    """Render a purposeful before/after comparison, not generic cards."""
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    gap = 0.46
    left = 0.76
    cw = (11.82 - gap) / 2
    max_ch = area.height - 0.15
    panels = [spec["left"], spec["right"]]
    tw = cw - 0.76
    def measure(size, bgap):
        cont = [sum(len(wrap_text(b, tw, size)) * line_height_in(size, 1.2) + bgap
                    for b in p["bullets"]) - bgap for p in panels]
        return cont, max(cont)

    def candidates():
        for bgap in stepped(0.32, 0.22, 0.02):
            _cont, used = measure(14, bgap)
            yield ("standard" if bgap == 0.32 else "gap",
                   {"size": 14, "gap": bgap}, used)
        for size in stepped(13.5, 11, 0.5):
            _cont, used = measure(size, 0.22)
            yield "font", {"size": size, "gap": 0.22}, used

    fitted = select_fit(
        "twocol", max_ch - 1.38, candidates(),
        guidance="左右の箇条書きを減らすか文言を短くしてください。",
    )
    size, bgap = fitted.values["size"], fitted.values["gap"]
    cont, _used = measure(size, bgap)
    body_h = max(cont) + 0.28
    top = area.top + 0.34
    for i, p in enumerate(panels):
        x = left + i * (cw + gap)
        fill = ZEBRA if i == 0 else LIGHT
        add_rect(slide, x, top, cw, min(max_ch - 0.14, body_h + 1.38), fill)
        add_text(slide, x + 0.36, top + 0.28, cw - 0.72, 0.25,
                 p.get("label", "BEFORE" if i == 0 else "AFTER"), 9.5,
                 bold=True, color=GRAY if i == 0 else ACCENT)
        heading_size, _ = fit_text_or_raise(
            "twocol", f"{'left' if i == 0 else 'right'}.heading",
            p["heading"], cw - 0.72, 0.48, 18,
            min_pt=14, weight="bold", spacing=1.1)
        add_text(slide, x + 0.36, top + 0.67, cw - 0.72, 0.48,
                 p["heading"], heading_size, bold=True, color=NAVY)
        y = top + 1.42
        for b in p["bullets"]:
            bh = len(wrap_text(b, tw, size)) * line_height_in(size, 1.2)
            add_text(slide, x + 0.36, y - 0.01, 0.24, 0.3, "—", 12,
                     bold=True, color=GRAY if i == 0 else ACCENT)
            add_text(slide, x + 0.7, y, tw, bh + 0.08, b, size,
                     color=TEXT, spacing=1.2)
            y += bh + bgap


def s_chart(slide, spec, page):
    area = header(slide, spec["kicker"], spec["title"], spec.get("lead"))
    categories = spec["chart"]["categories"]
    series = spec["chart"]["series"]
    kind = spec["chart"].get("kind", "bar")
    chart_types = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    }
    if kind not in chart_types:
        raise FitError(f"chart.kind={kind!r} は未対応です。")
    if not 1 <= len(categories) <= 12 or not 1 <= len(series) <= 4:
        raise FitError(
            "chart: 描画可能範囲はカテゴリ1〜12件、系列1〜4件です。"
            "カテゴリまたは系列を減らしてください。")
    fitted = _fit_chart_layout(
        area, len(categories), len(series), kind)
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(
        chart_types[kind], Inches(0.9), Inches(area.top + 0.22),
        Inches(11.8), Inches(area.height - 0.48), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = spec["chart"].get("show_legend", len(series) > 1)
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.font.size = Pt(fitted.values["font"])
    ch.font.name = FONT
    plot = ch.plots[0]
    show_values = spec["chart"].get("show_values", kind != "line")
    plot.has_data_labels = show_values
    if show_values:
        plot.data_labels.show_value = True
        plot.data_labels.font.size = Pt(fitted.values["label_font"])
        if spec["chart"].get("number_format"):
            plot.data_labels.number_format = spec["chart"]["number_format"]
            plot.data_labels.number_format_is_linked = False
    if hasattr(plot, "gap_width"):
        plot.gap_width = fitted.values["gap_width"]
    ch.value_axis.major_gridlines.format.line.color.rgb = RULE
    ch.value_axis.major_gridlines.format.line.width = Pt(0.6)
    ch.value_axis.format.line.fill.background()
    ch.category_axis.format.line.fill.background()
    for s, colr in zip(ch.series, (GRAY, ACCENT, CORAL, NAVY)):
        if kind == "line":
            s.format.line.color.rgb = colr
            s.format.line.width = Pt(2.0)
            s.marker.format.fill.solid()
            s.marker.format.fill.fore_color.rgb = colr
            s.marker.format.line.color.rgb = colr
        else:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colr
    if spec.get("note"):
        note_line(slide, spec["note"])


def _fit_chart_layout(area, category_count, series_count, kind):
    """軸ラベル間隔を圧縮してから、グラフ内文字と棒間隔を縮小する。"""
    legend = 0.34 if series_count > 1 else 0.0
    horizontal = kind in {"bar", "stacked_bar"}
    available = area.height - 0.48 if horizontal else 11.8
    unit = 1.0 if horizontal else 2.35
    candidates = [
        ("standard", {"font": 12.0, "label_font": 11.5, "gap_width": 76},
         category_count * 0.48 * unit + legend),
        ("gap", {"font": 11.0, "label_font": 10.5, "gap_width": 58},
         category_count * 0.40 * unit + legend),
        ("element", {"font": 9.5, "label_font": 9.0, "gap_width": 42},
         category_count * 0.32 * unit + legend),
    ]
    return select_fit(
        "chart", available, candidates,
        guidance="カテゴリ数を減らすか複数スライドへ分割してください。",
    )


RENDER = {"title": s_title, "bullets": s_bullets, "cards": s_cards,
          "table": s_table, "twocol": s_twocol, "chart": s_chart}


def render_slide(renderer, slide, spec, idx):
    """rendererの収容エラーをスライド位置つきの運用メッセージへ変換する。"""
    try:
        renderer(slide, spec, idx)
    except (ValueError, FileNotFoundError) as e:
        raise SystemExit(
            f"NG: slides[{idx - 1}] (type={spec['type']}) の生成に失敗:\n"
            f"  {e}") from e


def main(out_path, cover_footer_config=None):
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        configure_cover_footer(cover_footer_config)
    except ValueError as e:
        raise SystemExit(f"NG: 表紙・フッター設定: {e}") from e
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    total = len(DECK["slides"])
    for idx, spec in enumerate(DECK["slides"], 1):
        slide = prs.slides.add_slide(blank)
        render_slide(RENDER[spec["type"]], slide, spec, idx)
        if spec["type"] != "title":
            footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({total} slides)")


if __name__ == "__main__":
    default_out = Path(__file__).resolve().parent.parent / "out" / "sample_basic.pptx"
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out_path", nargs="?", default=default_out)
    parser.add_argument("--cover-footer-config", metavar="PATH",
                        help="表紙・フッター設定JSON")
    args = parser.parse_args()
    main(args.out_path, args.cover_footer_config)
