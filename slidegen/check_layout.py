"""生成済みPPTXのレイアウト衝突を機械検知する品質ゲート。

使い方: python slidegen/check_layout.py out/deck.pptx
検知対象:
  T-T: テキストグリフ同士の交差
  T-P: テキストグリフ×画像(アイコン)
  T-S: テキストグリフ×塗り図形の「部分重なり」(完全内包=意図的デザインは許可)
  T-F: テキストグリフがコンテナ枠線をまたぐ
  L-T: 矢印・線×テキストグリフ(白塗りマスクラベルは除外)
  CELL-OOB: 表セルからのテキストはみ出し
  OOB: スライド境界からの図形・画像・表・グラフのはみ出し
"""
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from cover_footer import COVER_BACKGROUND_NAME
from textfit import line_height_in, wrap_text

EMU = 914400
EDGE = 0.03          # 枠線の当たり判定幅
SEG_TRIM = 0.08      # 線分端はノード接続なので判定から除外する長さ
EPS = 0.03           # 視認できない接触(スリバー)を無視する許容量


def rect_of(sh):
    return (sh.left / EMU, sh.top / EMU, (sh.left + sh.width) / EMU,
            (sh.top + sh.height) / EMU)


def intersects(a, b, eps=EPS):
    return not (a[2] - eps <= b[0] or b[2] - eps <= a[0]
                or a[3] - eps <= b[1] or b[3] - eps <= a[1])


def contains(outer, inner, eps=EPS):
    return (outer[0] - eps <= inner[0] and outer[1] - eps <= inner[1]
            and outer[2] + eps >= inner[2] and outer[3] + eps >= inner[3])


def glyph_rect_for_text_frame(tf, box):
    """指定領域にあるtext frameの実グリフ矩形をtextfitで推定する。"""
    bw = box[2] - box[0]
    lines, sizes = [], []
    align = PP_ALIGN.LEFT
    for p in tf.paragraphs:
        if not p.runs:
            continue
        size = p.runs[0].font.size.pt if p.runs[0].font.size else 11.0
        bold = bool(p.runs[0].font.bold)
        text = "".join(r.text for r in p.runs)
        weight = "bold" if bold else "regular"
        for ln in wrap_text(text, bw, size, weight):
            lines.append((ln, size, weight))
            sizes.append(size)
        if p.alignment is not None:
            align = p.alignment
    if not lines:
        return None
    from textfit import text_width_in
    w = max(text_width_in(t, s, wt) for t, s, wt in lines)
    h = sum(line_height_in(s) for _, s, _ in lines)
    if align == PP_ALIGN.CENTER:
        x0 = box[0] + (bw - w) / 2
    elif align == PP_ALIGN.RIGHT:
        x0 = box[2] - w
    else:
        x0 = box[0]
    if tf.vertical_anchor == MSO_ANCHOR.MIDDLE:
        y0 = box[1] + ((box[3] - box[1]) - h) / 2
    else:
        y0 = box[1]
    return (x0, y0, x0 + w, y0 + h)


def glyph_rect(sh):
    """通常図形の実グリフ矩形をtextfitで推定する。"""
    return glyph_rect_for_text_frame(sh.text_frame, rect_of(sh))


def table_cells(sh):
    """表セルの外形・文字領域とtext frameを返す。"""
    table = sh.table
    x0, y0, _, _ = rect_of(sh)
    col_widths = [col.width / EMU for col in table.columns]
    row_heights = [row.height / EMU for row in table.rows]
    y = y0
    for ri, row in enumerate(table.rows):
        x = x0
        row_h = row_heights[ri]
        for ci, col in enumerate(table.columns):
            col_w = col_widths[ci]
            cell = table.cell(ri, ci)
            if not getattr(cell, "is_spanned", False):
                span_w = getattr(cell, "span_width", 1)
                span_h = getattr(cell, "span_height", 1)
                merged_w = sum(col_widths[ci:ci + span_w])
                merged_h = sum(row_heights[ri:ri + span_h])
                outer = (x, y, x + merged_w, y + merged_h)
                inner = (
                    x + cell.margin_left / EMU,
                    y + cell.margin_top / EMU,
                    x + merged_w - cell.margin_right / EMU,
                    y + merged_h - cell.margin_bottom / EMU,
                )
                yield outer, inner, cell.text_frame, ri, ci
            x += col_w
        y += row_h


def is_oob(rect, slide_w, slide_h):
    return (rect[0] < -EPS or rect[1] < -EPS
            or rect[2] > slide_w + EPS or rect[3] > slide_h + EPS)


def seg_of(sh):
    """コネクタの線分端点(flipで向きを解決)。"""
    x1, y1, x2, y2 = rect_of(sh)
    fh = getattr(sh, "rotation", 0)  # dummy no-op
    el = sh._element
    flip_h = el.spPr.xfrm.get("flipH") == "1" if el.spPr.xfrm is not None else False
    flip_v = el.spPr.xfrm.get("flipV") == "1" if el.spPr.xfrm is not None else False
    ax, bx = (x2, x1) if flip_h else (x1, x2)
    ay, by = (y2, y1) if flip_v else (y1, y2)
    return (ax, ay, bx, by)


def seg_hits_rect(seg, r):
    """線分と矩形の交差(端をSEG_TRIMだけ縮めて接続点は無視)。"""
    x1, y1, x2, y2 = seg
    import math
    L = math.hypot(x2 - x1, y2 - y1)
    if L < SEG_TRIM * 2:
        return False
    ux, uy = (x2 - x1) / L, (y2 - y1) / L
    x1, y1 = x1 + ux * SEG_TRIM, y1 + uy * SEG_TRIM
    x2, y2 = x2 - ux * SEG_TRIM, y2 - uy * SEG_TRIM
    # 双方の投影が重なるかをサンプリングで簡易判定
    steps = max(2, int(L / 0.05))
    for i in range(steps + 1):
        t = i / steps
        px, py = x1 + (x2 - x1) * t, y1 + (y2 - y1) * t
        if r[0] + EPS < px < r[2] - EPS and r[1] + EPS < py < r[3] - EPS:
            return True
    return False


def has_solid_fill(sh):
    try:
        return sh.fill.type is not None and str(sh.fill.type) == "SOLID (1)"
    except Exception:
        return False


def snippet(t):
    t = t.replace("\n", " ")
    return t[:16] + ("…" if len(t) > 16 else "")


def check(path):
    prs = Presentation(path)
    slide_w = prs.slide_width / EMU
    slide_h = prs.slide_height / EMU
    findings = []
    for si, slide in enumerate(prs.slides, 1):
        texts, pics, solids, frames, segs = [], [], [], [], []
        for z, sh in enumerate(slide.shapes):  # zは描画順(後勝ち)
            st = sh.shape_type
            bounds = rect_of(sh)
            if is_oob(bounds, slide_w, slide_h):
                findings.append((si, "OOB", sh.name, ""))
            if st == MSO_SHAPE_TYPE.PICTURE:
                if sh.name != COVER_BACKGROUND_NAME:
                    pics.append((bounds, sh.name))
            elif st in (MSO_SHAPE_TYPE.AUTO_SHAPE,):
                (solids if has_solid_fill(sh) else frames).append(
                    (bounds, sh.name, z))
            elif st == MSO_SHAPE_TYPE.LINE:
                segs.append((seg_of(sh), sh.name, z))
            elif st == MSO_SHAPE_TYPE.CHART:
                pics.append((bounds, sh.name))
            elif st == MSO_SHAPE_TYPE.TABLE:
                for _cell_rect, text_rect, tf, ri, ci in table_cells(sh):
                    if not tf.text.strip():
                        continue
                    g = glyph_rect_for_text_frame(tf, text_rect)
                    if g:
                        label = snippet(tf.text)
                        texts.append((g, label, False))
                        if not contains(text_rect, g, eps=0.01):
                            findings.append(
                                (si, "CELL-OOB", f"{sh.name}[{ri},{ci}]", label))
            if st != MSO_SHAPE_TYPE.TABLE \
                    and sh.has_text_frame and sh.text_frame.text.strip():
                g = glyph_rect(sh)
                if g:
                    texts.append((g, snippet(sh.text_frame.text),
                                  has_solid_fill(sh)))
        # T-T
        for i in range(len(texts)):
            for j in range(i + 1, len(texts)):
                if intersects(texts[i][0], texts[j][0]):
                    findings.append((si, "T-T", texts[i][1], texts[j][1]))
        # T-P
        for g, t, _ in texts:
            for r, name in pics:
                if intersects(g, r):
                    findings.append((si, "T-P", t, name))
        # T-S: 部分重なりのみ(内包は許可)
        for g, t, _ in texts:
            for r, name, _z in solids:
                if intersects(g, r) and not contains(r, g):
                    findings.append((si, "T-S", t, name))
        # T-F: 枠線またぎ(内包/完全外は許可。白塗りマスクラベルは枠線を隠すので許可)
        for g, t, masked in texts:
            if masked:
                continue
            for r, name, _z in frames:
                if intersects(g, r, eps=-EDGE) and not contains(
                        (r[0] + EDGE, r[1] + EDGE, r[2] - EDGE, r[3] - EDGE), g) \
                        and intersects(g, r):
                    findings.append((si, "T-F", t, name))
        # L-T (マスクラベルと、線より後に描かれた塗り図形上のテキストは除外)
        for seg, name, lz in segs:
            for g, t, masked in texts:
                if masked or not seg_hits_rect(seg, g):
                    continue
                covered = any(contains(r, g) and sz > lz
                              for r, _n, sz in solids)
                if not covered:
                    findings.append((si, "L-T", name, t))
        # L-P
        for seg, name, _lz in segs:
            for r, pname in pics:
                if seg_hits_rect(seg, r):
                    findings.append((si, "L-P", name, pname))
        for g, t, _ in texts:
            if is_oob(g, slide_w, slide_h):
                findings.append((si, "OOB-TEXT", t, ""))
    return findings


if __name__ == "__main__":
    path = sys.argv[1]
    fs = check(path)
    if not fs:
        print(f"OK: no layout collisions in {path}")
        sys.exit(0)
    print(f"NG: {len(fs)} finding(s) in {path}")
    for si, kind, a, b in fs:
        print(f"  slide{si:02d} [{kind}] {a!r} x {b!r}")
    sys.exit(1)
