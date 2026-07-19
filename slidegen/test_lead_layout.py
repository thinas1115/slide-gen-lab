"""任意leadと本文領域縮小の回帰テスト。"""
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches

import generate
from content_lead_patterns import LEAD_PATTERN_DECK, LEADS
from generate_from_json import RENDER
from layout_fit import FitError
from validate_content import validate


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    return prs.slides.add_slide(prs.slide_layouts[6])


def main():
    default_area = generate.header(_slide(), "分類", "タイトル")
    assert default_area.top == generate.BODY_TOP
    assert not default_area.shifted

    lead_area = generate.header(
        _slide(), "分類", "タイトル", "結論を先に示してから本文を説明します。")
    assert lead_area.top > generate.BODY_TOP
    assert lead_area.bottom == generate.BODY_BOTTOM
    assert lead_area.shifted

    errors = validate(deepcopy(LEAD_PATTERN_DECK))
    assert not errors, "\n".join(errors)

    generate.DECK = LEAD_PATTERN_DECK
    for idx, spec in enumerate(LEAD_PATTERN_DECK["slides"], 1):
        generate.render_slide(RENDER[spec["type"]], _slide(), spec, idx)

    invalid = deepcopy(LEAD_PATTERN_DECK)
    invalid["slides"][0]["lead"] = ""
    assert any("空でない文字列" in e
               for e in validate(invalid))

    title_with_lead = {
        "meta": deepcopy(LEAD_PATTERN_DECK["meta"]),
        "slides": [{"type": "title", "title": "表紙", "subtitle": "副題",
                    "lead": "表紙では使わない"}],
    }
    assert any("表紙以外" in e for e in validate(title_with_lead))

    try:
        generate.header(_slide(), "分類", "タイトル", "長いリード文" * 200)
    except FitError as e:
        assert "header.lead" in str(e), str(e)
        assert "文言を短く" in str(e), str(e)
    else:
        raise AssertionError("過密なleadを拒否しませんでした")

    dense = deepcopy(LEAD_PATTERN_DECK["slides"][-1])
    dense["lead"] = LEADS["diagram"]
    dense["diagram"]["rows"].append("overflow")
    overflow = deepcopy(next(iter(dense["diagram"]["nodes"].values())))
    overflow.update(row="overflow", title="追加ノード", sub="追加説明")
    dense["diagram"]["nodes"]["overflow"] = overflow
    try:
        RENDER["diagram"](_slide(), dense, 1)
    except FitError as e:
        assert "アイコン縮小" in str(e), str(e)
        assert "行数を減らす" in str(e), str(e)
    else:
        raise AssertionError("2行leadと追加行で過密になるdiagramを拒否しませんでした")

    print(f"lead layout tests passed ({len(LEAD_PATTERN_DECK['slides'])} patterns)")


if __name__ == "__main__":
    main()
