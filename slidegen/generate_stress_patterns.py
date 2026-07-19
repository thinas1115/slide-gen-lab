"""段階的縮小ストレス検証デッキを生成する。"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import generate
from content_stress_patterns import STRESS_PATTERN_DECK
from generate_from_json import RENDER
from validate_content import validate


def main(out_path, cover_footer_config=None):
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        generate.configure_cover_footer(cover_footer_config)
    except ValueError as e:
        raise SystemExit(f"NG: 表紙・フッター設定: {e}") from e
    errors = validate(STRESS_PATTERN_DECK, allow_sample_content=True)
    if errors:
        raise SystemExit("NG: ストレスギャラリー\n  - " + "\n  - ".join(errors))
    generate.DECK = STRESS_PATTERN_DECK
    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(STRESS_PATTERN_DECK["slides"], 1):
        slide = prs.slides.add_slide(blank)
        generate.render_slide(RENDER[spec["type"]], slide, spec, idx)
        generate.footer(slide, idx)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(STRESS_PATTERN_DECK['slides'])} slides)")


if __name__ == "__main__":
    default_out = Path(__file__).resolve().parent.parent / "out" / "stress_gallery.pptx"
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out_path", nargs="?", default=default_out)
    parser.add_argument("--cover-footer-config", metavar="PATH",
                        help="表紙・フッター設定JSON")
    args = parser.parse_args()
    main(args.out_path, args.cover_footer_config)
