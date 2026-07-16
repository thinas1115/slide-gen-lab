"""content.json を入力にしてPPTXを生成する。

使い方:
  python slidegen/generate_from_json.py content.json out/from_json.pptx

生成前に validate_content.py の検証を通す。エラーがあればレンダリングせず、
生成AIにそのまま渡して直させられる粒度のメッセージを出して終了する。
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import generate
from diagrams import s_hub, s_org
from diagrams2 import s_matrix, s_process, s_program_roadmap, s_roadmap
from diagram_layout import render_diagram
from image_slide import s_image
from validate_content import validate


def s_diagram(slide, spec, page):
    """構成図: グリッド仕様(spec["diagram"])から座標ゼロで自動レイアウト。

    仕様はスライド内にインラインで書く。名前付きテンプレート参照は持たない。
    """
    area = generate.header(
        slide, spec["kicker"], spec["title"], spec.get("lead"))
    render_diagram(slide, spec["diagram"], note=spec.get("note"),
                   content_area=area if area.shifted else None)


# 旧固定構成図typeの aws / aws2 は廃止済み。任意テーマの構成図は
# diagram type(宣言的レイアウトエンジン)で生成する。
RENDER = dict(generate.RENDER,
              hub=s_hub, org=s_org,
              process=s_process, roadmap=s_roadmap,
              program_roadmap=s_program_roadmap, matrix=s_matrix,
              diagram=s_diagram, image=s_image)


def main(json_path, out_path, cover_footer_config=None):
    deck = json.loads(Path(json_path).read_text(encoding="utf-8"))
    errors = validate(deck)
    if errors:
        print(f"NG: {json_path} に {len(errors)} 件の問題 (生成を中止):",
              file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        raise SystemExit(1)
    generate.DECK = deck
    try:
        generate.configure_cover_footer(cover_footer_config)
    except ValueError as e:
        raise SystemExit(f"NG: 表紙・フッター設定: {e}") from e

    prs = Presentation()
    prs.slide_width = Inches(generate.SLIDE_W)
    prs.slide_height = Inches(generate.SLIDE_H)
    blank = prs.slide_layouts[6]
    for idx, spec in enumerate(deck["slides"], 1):
        slide = prs.slides.add_slide(blank)
        try:
            renderer = RENDER[spec["type"]]
        except KeyError as e:
            known = ", ".join(sorted(RENDER))
            raise SystemExit(f"unknown slide type: {spec['type']!r}. known: {known}") from e
        generate.render_slide(renderer, slide, spec, idx)
        if spec["type"] != "title":
            generate.footer(slide, idx)
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"saved: {out_path} ({len(deck['slides'])} slides)")


if __name__ == "__main__":
    project_root = Path(__file__).resolve().parent.parent
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("json_path", nargs="?", default=project_root / "content.json")
    parser.add_argument("out_path", nargs="?", default=project_root / "out" / "from_json.pptx")
    parser.add_argument("--cover-footer-config", metavar="PATH",
                        help="表紙・フッター設定JSON")
    args = parser.parse_args()
    main(args.json_path, args.out_path, args.cover_footer_config)
