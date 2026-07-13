"""AWSアイコンデッキから対象サービスのアイコンPNGを抽出する。

シェイプ名にサービス名がないため、「ラベルテキストの直上にある最寄りの
Picture」を対応アイコンとみなして抽出する。
"""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SRC = r"<AWS_ICON_DECK_PATH>"
DEST = Path(__file__).parent / "assets"
DEST.mkdir(exist_ok=True)

# 出力名: (完全一致優先の候補ラベル)
TARGETS = {
    "bedrock": ["Amazon Bedrock"],
    "s3": ["Amazon Simple Storage Service (Amazon S3)", "Amazon Simple Storage Service"],
    "fargate": ["AWS Fargate", "Amazon Elastic Container Service"],
    "cloudwatch": ["Amazon CloudWatch"],
    "alb": ["Elastic Load Balancing", "Application Load Balancer"],
    "users": ["Users"],
    "user": ["User"],
    "rds": ["Amazon Relational Database Service (Amazon RDS)"],
    "cloudfront": ["Amazon CloudFront"],
    "route53": ["Amazon Route 53"],
    "dynamodb": ["Amazon DynamoDB"],
    "sqs": ["Amazon Simple Queue Service (Amazon SQS)"],
    "ecr": ["Amazon Elastic Container Registry (Amazon ECR)"],
}


def walk(shapes, ox=0, oy=0):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 子座標系→親座標系の変換(スケールは1前提の近似)
            cx = sh.left - sh.shapes[0].left if False else 0
            yield from walk(sh.shapes, ox + sh.left, oy + sh.top)
        else:
            try:
                yield sh, ox + (sh.left or 0), oy + (sh.top or 0)
            except TypeError:
                continue


prs = Presentation(SRC)
found = {}
for si, slide in enumerate(prs.slides):
    pics, texts = [], []
    for sh, ax, ay in walk(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics.append((sh, ax, ay))
        elif sh.has_text_frame:
            t = sh.text_frame.text.strip().replace(" ", " ").replace("\n", " ")
            if t:
                texts.append((t, ax, ay, sh.width or 0))
    if not pics:
        continue
    for key, labels in TARGETS.items():
        if key in found:
            continue
        for label in labels:
            cand = [(t, x, y, w) for t, x, y, w in texts if t.lower() == label.lower()]
            if not cand:
                continue
            t, tx, ty, tw = cand[0]
            tcx = tx + tw / 2
            best, bd = None, None
            for sh, px, py in pics:
                if py > ty:          # アイコンはラベルより上にある
                    continue
                pcx = px + (sh.width or 0) / 2
                d = abs(pcx - tcx) + (ty - py) * 0.3
                if bd is None or d < bd:
                    best, bd = sh, d
            if best is not None:
                img = best.image
                ext = img.ext
                p = DEST / f"{key}.{ext}"
                p.write_bytes(img.blob)
                found[key] = f"slide{si + 1} label={t!r} -> {p.name} ({len(img.blob)}B, {ext})"
                break

report = [f"{k}: {v}" for k, v in found.items()]
missing = [k for k in TARGETS if k not in found]
report.append(f"missing: {missing}")
Path(__file__).with_name("aws_icons_report.txt").write_text(
    "\n".join(report), encoding="utf-8")
print("done", len(found), "found /", len(missing), "missing")
